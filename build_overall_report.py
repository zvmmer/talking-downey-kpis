"""Talking Downey — cumulative all-time performance report.

Rolls up EVERY tracked post since engagement start (2026-06-08) into a
single deck. This is the "since we started" view — no monthly split,
no comparison — just the full picture for someone asking "how have we
done since Kaname took over?"

Uses report_helpers.py.

Usage
-----
    MCP/.venv/bin/python Projects/talking_downey/kpis/build_overall_report.py
"""
from __future__ import annotations

import argparse
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from report_helpers import (
    PALETTE, LAYOUT_W, LAYOUT_H, ENGAGEMENT_START,
    add_rect, add_text, fmt_int,
    load_all_snapshots, posts_in_window, group_by_category,
    slide_new, slide_header, slide_footer_takeaway,
)

TODAY = date.today().isoformat()


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    add_rect(slide, 0, 0, 0.45, LAYOUT_H, PALETTE["cherry"])

    add_text(slide, 1.0, 1.2, 11, 0.5,
             "OVERALL PERFORMANCE — ALL-TIME",
             size=14, color=PALETTE["cherry"], bold=True)
    add_text(slide, 1.0, 1.9, 11, 1.3,
             "Talking Downey", size=64, color=PALETTE["white"],
             bold=True, font="Georgia")
    add_text(slide, 1.0, 3.4, 11, 0.6,
             f"Since Jun 8, 2026 · {TODAY}",
             size=24, color=PALETTE["cream"], font="Georgia")

    add_text(slide, 1.0, 5.0, 8, 0.4,
             "Every tracked post across every platform.",
             size=13, color=PALETTE["muted"])
    add_text(slide, 1.0, 5.4, 8, 0.4,
             "The cumulative story — no monthly cutoffs.",
             size=13, color=PALETTE["muted"])

    add_text(slide, 1.0, 6.6, 11, 0.4,
             "Prepared by Zhamir Pascual — Kaname Z",
             size=12, color=PALETTE["muted"])


def slide_totals(prs, all_posts):
    slide = slide_new(prs)
    slide_header(slide, "The headline",
                 "Total reach since engagement started",
                 "Sum of view counts across every unique post captured since Jun 8, 2026.")

    total_views = sum(int(p.get("view_count") or 0) for p in all_posts)
    total_posts = len(all_posts)
    total_likes = sum(int(p.get("like_count") or 0) for p in all_posts)

    add_text(slide, 0.8, 2.4, 11.7, 1.6,
             fmt_int(total_views), size=110, color=PALETTE["cherry"],
             bold=True, font="Georgia", align="center")
    add_text(slide, 0.8, 4.0, 11.7, 0.5,
             "total views across all tracked posts",
             size=18, color=PALETTE["navy"], align="center")

    card_w, gap = 3.8, 0.3
    total_w = card_w * 3 + gap * 2
    x_start = (LAYOUT_W - total_w) / 2
    y = 5.0
    for i, (label, val) in enumerate([
        ("Total posts", fmt_int(total_posts)),
        ("Total likes", fmt_int(total_likes)),
        ("Avg views/post", fmt_int(total_views // total_posts if total_posts else 0)),
    ]):
        x = x_start + i * (card_w + gap)
        add_rect(slide, x, y, card_w, 1.4, PALETTE["white"])
        add_rect(slide, x, y, 0.08, 1.4, PALETTE["cherry"])
        add_text(slide, x + 0.3, y + 0.2, card_w - 0.5, 0.3,
                 label.upper(), size=11, color=PALETTE["cherry"], bold=True)
        add_text(slide, x + 0.3, y + 0.55, card_w - 0.5, 0.8,
                 val, size=32, color=PALETTE["ink"], bold=True, font="Georgia")


def slide_platforms(prs, all_posts):
    slide = slide_new(prs)
    slide_header(slide, "By platform",
                 "Where the reach is coming from",
                 "Total views + post count per platform since Jun 8. Bigger view total = the audience is there.")

    platforms = [("instagram", "Instagram", PALETTE["cherry"]),
                 ("tiktok",    "TikTok",    PALETTE["ink"]),
                 ("facebook",  "Facebook",  PALETTE["navy"]),
                 ("youtube",   "YouTube",   PALETTE["cherry"])]
    card_w, gap = 2.85, 0.2
    total_w = card_w * 4 + gap * 3
    x_start = (LAYOUT_W - total_w) / 2
    y = 2.4

    for i, (key, label, color) in enumerate(platforms):
        x = x_start + i * (card_w + gap)
        plat_posts = [p for p in all_posts if p.get("platform") == key]
        views = sum(int(p.get("view_count") or 0) for p in plat_posts)
        n = len(plat_posts)

        add_rect(slide, x, y, card_w, 3.8, PALETTE["white"])
        add_rect(slide, x, y, 0.06, 3.8, color)
        add_text(slide, x + 0.2, y + 0.2, card_w - 0.4, 0.4,
                 label.upper(), size=11, color=color, bold=True)
        add_text(slide, x + 0.2, y + 0.65, card_w - 0.4, 1.1,
                 fmt_int(views), size=42, color=PALETTE["navy"],
                 bold=True, font="Georgia")
        add_text(slide, x + 0.2, y + 1.8, card_w - 0.4, 0.3,
                 "total views", size=11, color=PALETTE["muted"])
        add_text(slide, x + 0.2, y + 2.3, card_w - 0.4, 0.3,
                 f"{n} posts", size=13, color=PALETTE["ink"], bold=True)
        avg = views / n if n else 0
        add_text(slide, x + 0.2, y + 2.7, card_w - 0.4, 0.3,
                 f"{avg:,.0f} avg views/post",
                 size=11, color=PALETTE["ink"])

    slide_footer_takeaway(slide,
        "Instagram is the workhorse. TikTok punches above weight on breakouts. YouTube compounds. FB is the reliable floor.")


def slide_category_ranking(prs, all_posts):
    slide = slide_new(prs)
    slide_header(slide, "All-time category ranking",
                 "Which content types earned the most reach",
                 "Cumulative view totals per category. Numbers speak — this is the sponsor's decision matrix.")

    buckets = group_by_category(all_posts)
    scored = []
    for cat, posts in buckets.items():
        if not posts:
            continue
        views = sum(int(p.get("view_count") or 0) for p in posts)
        scored.append({"cat": cat, "posts": len(posts), "views": views,
                       "avg": views / len(posts) if posts else 0})
    scored.sort(key=lambda s: -s["views"])
    top = scored[:12]

    max_views = max((s["views"] for s in top), default=1)

    label_map = {
        "trujillo_scandal":              ("Trujillo Scandal",       "cherry"),
        "lisette_scandal":               ("Lisette Scandal",        "cherry"),
        "ai_fake_news":                  ("AI Fake News",           "cherry"),
        "council_chaos":                 ("Council Chaos",          "cherry"),
        "ice_immigration":               ("ICE / Immigration",      "cherry"),
        "elections":                     ("Elections",              "navy"),
        "judicial":                      ("Judicial",               "navy"),
        "politics_other":                ("Politics (other)",       "navy"),
        "political_individual_highlight":("Political Highlight",    "sky"),
        "resident_highlight":            ("Resident Highlight",     "green"),
        "downtown_development":          ("Downtown Development",   "gold"),
        "american_pride":                ("American Pride",         "cherry"),
        "fifa_positive":                 ("FIFA — Positive",        "green"),
        "fifa_critique":                 ("FIFA — Critique",        "muted"),
        "food":                          ("Food",                   "green"),
        "brand_growth":                  ("Brand Growth",           "blush"),
        "community_events":              ("Community Events",       "ink"),
        "community_local":               ("Community Local",        "ink"),
        "other":                         ("Other",                  "muted"),
    }

    row_h = 0.36
    y0 = 2.35
    x_bar = 3.7
    bar_max_w = 7.5

    for i, s in enumerate(top):
        y = y0 + i * row_h
        label, color_key = label_map.get(s["cat"], (s["cat"], "muted"))
        add_text(slide, 0.8, y + 0.03, 2.7, row_h,
                 label, size=11, color=PALETTE["ink"], bold=(i < 3))
        bar_w = max(0.05, bar_max_w * s["views"] / max_views)
        add_rect(slide, x_bar, y + 0.06, bar_w, row_h - 0.15, PALETTE[color_key])
        add_text(slide, x_bar + bar_w + 0.1, y + 0.03, 2.5, row_h,
                 f"{fmt_int(s['views'])}   ({s['posts']} posts · {s['avg']:,.0f} avg)",
                 size=10, color=PALETTE["muted"])

    slide_footer_takeaway(slide,
        "The top of the leaderboard tells you what's earning attention. The tail tells you where headroom exists.")


def slide_top10(prs, all_posts):
    slide = slide_new(prs)
    slide_header(slide, "All-time top 10 posts",
                 "The peak performers since Jun 8",
                 "Individual URLs ranked by views. These are the moments the audience shows up hardest.")

    posts = sorted(all_posts, key=lambda p: -int(p.get("view_count") or 0))[:10]

    col_widths = [0.5, 1.0, 6.4, 1.4, 1.4, 1.4]
    headers = ["#", "PLATFORM", "POST", "VIEWS", "LIKES", "CATEGORY"]
    row_y = 2.4
    row_h = 0.42
    x_base = 0.6

    add_rect(slide, x_base, row_y, sum(col_widths), row_h, PALETTE["navy"])
    cx = x_base
    for hd, w in zip(headers, col_widths):
        add_text(slide, cx + 0.1, row_y + 0.08, w - 0.2, 0.3,
                 hd, size=10, color=PALETTE["cream"], bold=True)
        cx += w

    for i, p in enumerate(posts, 1):
        ry = row_y + row_h + (i - 1) * row_h
        if i % 2 == 0:
            add_rect(slide, x_base, ry, sum(col_widths), row_h, PALETTE["soft"])
        if i <= 3:
            add_rect(slide, x_base, ry, 0.06, row_h, PALETTE["cherry"])

        title = (p.get("title") or p.get("description_excerpt") or "")[:80]
        plat = {"instagram": "IG", "tiktok": "TT", "facebook": "FB", "youtube": "YT"}.get(p.get("platform"), "?")
        cat = p.get("_category", "other").replace("_", " ").title()

        cx = x_base
        add_text(slide, cx + 0.1, ry + 0.08, col_widths[0] - 0.2, 0.3,
                 str(i), size=12, color=PALETTE["cherry"], bold=True)
        cx += col_widths[0]
        add_text(slide, cx + 0.1, ry + 0.08, col_widths[1] - 0.2, 0.3,
                 plat, size=11, color=PALETTE["navy"], bold=True)
        cx += col_widths[1]
        add_text(slide, cx + 0.1, ry + 0.08, col_widths[2] - 0.2, 0.3,
                 title, size=10, color=PALETTE["ink"])
        cx += col_widths[2]
        add_text(slide, cx + 0.1, ry + 0.08, col_widths[3] - 0.2, 0.3,
                 fmt_int(p.get("view_count")), size=12, color=PALETTE["navy"], bold=True)
        cx += col_widths[3]
        add_text(slide, cx + 0.1, ry + 0.08, col_widths[4] - 0.2, 0.3,
                 fmt_int(p.get("like_count")), size=11, color=PALETTE["ink"])
        cx += col_widths[4]
        add_text(slide, cx + 0.1, ry + 0.08, col_widths[5] - 0.2, 0.3,
                 cat[:20], size=9, color=PALETTE["muted"])


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path,
                    default=Path(__file__).parent / "OVERALL_REPORT.pptx")
    ap.add_argument("--kpis-root", type=Path,
                    default=Path(__file__).parent)
    args = ap.parse_args()

    snaps = load_all_snapshots(args.kpis_root / "episodes_public")
    if not snaps:
        sys.exit(f"No snapshots found under {args.kpis_root/'episodes_public'}")

    all_posts = posts_in_window(snaps, ENGAGEMENT_START, None)
    print(f"Loaded {len(snaps)} snapshots · {len(all_posts)} unique posts since {ENGAGEMENT_START}")

    prs = Presentation()
    prs.slide_width = Inches(LAYOUT_W)
    prs.slide_height = Inches(LAYOUT_H)

    slide_title(prs)
    slide_totals(prs, all_posts)
    slide_platforms(prs, all_posts)
    slide_category_ranking(prs, all_posts)
    slide_top10(prs, all_posts)

    prs.save(args.out)
    print(f"✓ Wrote: {args.out}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
