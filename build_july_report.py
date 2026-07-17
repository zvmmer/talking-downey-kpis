"""Talking Downey — July 2026 Performance Report deck.

Cousin of build_june_report.py — same visual language, different spine.

The June report leads with the political engine narrative
("controversy drives reach"). July shifts the frame to
BRAND MOMENTUM — the audience widening, positive stories
(Downey growth, American pride) picking up alongside the political
work. The sponsor sees a portfolio, not a one-note show.

Uses report_helpers.py (shared with overall + sponsor decks).
The June builder stays self-contained by design — frozen baseline.

Usage
-----
    MCP/.venv/bin/python Projects/talking_downey/kpis/build_july_report.py

Outputs `MONTHLY_REPORT_2026-07.pptx` at the kpis/ root by default.
"""
from __future__ import annotations

import argparse
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from report_helpers import (
    PALETTE, LAYOUT_W, LAYOUT_H, ENGAGEMENT_START,
    add_rect, add_text, fmt_int, normalize_date,
    load_all_snapshots, posts_in_window, group_by_category,
    slide_new, slide_header, slide_footer_takeaway,
    CATEGORY_ORDER,
)

JULY_START = "2026-07-01"
JULY_END   = "2026-07-31"
JUNE_START = "2026-06-08"
JUNE_END   = "2026-06-30"


# ── slides ──────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    add_rect(slide, 0, 0, 0.45, LAYOUT_H, PALETTE["sky"])   # sky = growth spine

    add_text(slide, 1.0, 1.2, 11, 0.5,
             "JULY 2026 PERFORMANCE REPORT",
             size=14, color=PALETTE["sky"], bold=True)
    add_text(slide, 1.0, 1.9, 11, 1.3,
             "Talking Downey", size=64, color=PALETTE["white"],
             bold=True, font="Georgia")
    add_text(slide, 1.0, 3.4, 11, 0.6,
             "Brand momentum — widening the portfolio",
             size=24, color=PALETTE["cream"], font="Georgia")

    add_text(slide, 1.0, 5.0, 8, 0.4,
             "Tracking period: Jul 1 – Jul 31, 2026",
             size=13, color=PALETTE["muted"])
    add_text(slide, 1.0, 5.4, 8, 0.4,
             "Comparison baseline: June 2026 (16-day launch window)",
             size=13, color=PALETTE["muted"])
    add_text(slide, 1.0, 5.8, 8, 0.4,
             "4 platforms · IG · TikTok · Facebook · YouTube",
             size=13, color=PALETTE["muted"])

    add_text(slide, 1.0, 6.6, 11, 0.4,
             "Prepared by Zhamir Pascual — Kaname Z",
             size=12, color=PALETTE["muted"])


def slide_month_shift(prs, july_posts, june_posts):
    """The spine slide — 'here's what changed from June to July'."""
    slide = slide_new(prs)
    slide_header(slide, "June → July",
                 "The month the show broadened",
                 "Same brand. Wider net. Political content still runs — but Downey growth stories and American pride content joined the top tier.")

    # Big number: post volume ratio
    v_june = len(june_posts)
    v_july = len(july_posts)
    add_text(slide, 0.8, 2.4, 5.8, 1.5,
             fmt_int(v_july), size=110, color=PALETTE["sky"], bold=True,
             font="Georgia", align="center")
    add_text(slide, 0.8, 4.0, 5.8, 0.5,
             "posts in July", size=16, color=PALETTE["navy"], align="center")
    add_text(slide, 0.8, 4.5, 5.8, 0.4,
             f"({v_june} in June baseline window)",
             size=11, color=PALETTE["muted"], align="center")

    # Category shift snapshot
    def _view_avg(posts, cat_group):
        rel = [int(p.get("view_count") or 0) for p in posts if p.get("_category") in cat_group]
        return sum(rel) / len(rel) if rel else 0

    scandal_grp = ["trujillo_scandal", "lisette_scandal", "council_chaos", "ai_fake_news"]
    highlights_grp = ["resident_highlight", "political_individual_highlight"]
    american_grp = ["american_pride"]

    shifts = [
        ("Truth & Scandal",        _view_avg(june_posts, scandal_grp),    _view_avg(july_posts, scandal_grp),    PALETTE["cherry"]),
        ("Personality Highlights", _view_avg(june_posts, highlights_grp), _view_avg(july_posts, highlights_grp), PALETTE["sky"]),
        ("American Pride",         _view_avg(june_posts, american_grp),   _view_avg(july_posts, american_grp),   PALETTE["gold"]),
    ]

    y = 2.4
    x = 7.0
    for label, jun, jul, color in shifts:
        add_rect(slide, x, y, 5.6, 1.2, PALETTE["white"])
        add_rect(slide, x, y, 0.08, 1.2, color)
        add_text(slide, x + 0.25, y + 0.15, 4, 0.3,
                 label.upper(), size=10, color=color, bold=True)
        add_text(slide, x + 0.25, y + 0.45, 3, 0.6,
                 f"{jun:,.0f} → {jul:,.0f}",
                 size=22, color=PALETTE["ink"], bold=True, font="Georgia")
        arrow = "▲" if jul > jun else ("▼" if jul < jun else "▬")
        arrow_color = PALETTE["green"] if jul > jun else (PALETTE["cherry"] if jul < jun else PALETTE["muted"])
        pct = f"{((jul - jun) / jun * 100):+.0f}%" if jun else "—"
        add_text(slide, x + 4.2, y + 0.45, 1.3, 0.6,
                 f"{arrow} {pct}", size=18, color=arrow_color, bold=True)
        add_text(slide, x + 0.25, y + 0.85, 5, 0.3,
                 "avg views per post",
                 size=10, color=PALETTE["muted"])
        y += 1.4

    slide_footer_takeaway(slide,
        "The show is no longer a one-note engine. Positive lanes are earning their seat at the table.")


def slide_platform_growth(prs, snaps):
    """Total views per platform, June cumulative vs July cumulative."""
    slide = slide_new(prs)
    slide_header(slide, "Platform-by-platform",
                 "Where the audience is showing up",
                 "Total views per platform in the July window vs. the June baseline window. Positive = platform is still climbing.")

    june_posts = posts_in_window(snaps, JUNE_START, JUNE_END)
    july_posts = posts_in_window(snaps, JULY_START, JULY_END)

    platforms = [("instagram", "Instagram"), ("tiktok", "TikTok"),
                 ("facebook", "Facebook"), ("youtube", "YouTube")]
    card_w = 2.85
    gap = 0.2
    total_w = card_w * 4 + gap * 3
    x_start = (LAYOUT_W - total_w) / 2
    y = 2.5

    for i, (key, label) in enumerate(platforms):
        x = x_start + i * (card_w + gap)
        june_v = sum(int(p.get("view_count") or 0) for p in june_posts if p.get("platform") == key)
        july_v = sum(int(p.get("view_count") or 0) for p in july_posts if p.get("platform") == key)
        june_n = sum(1 for p in june_posts if p.get("platform") == key)
        july_n = sum(1 for p in july_posts if p.get("platform") == key)

        add_rect(slide, x, y, card_w, 3.8, PALETTE["white"])
        add_rect(slide, x, y, 0.06, 3.8, PALETTE["sky"])
        add_text(slide, x + 0.2, y + 0.2, card_w - 0.4, 0.4,
                 label.upper(), size=11, color=PALETTE["sky"], bold=True)

        add_text(slide, x + 0.2, y + 0.65, card_w - 0.4, 0.9,
                 fmt_int(july_v), size=36, color=PALETTE["navy"],
                 bold=True, font="Georgia")
        add_text(slide, x + 0.2, y + 1.55, card_w - 0.4, 0.3,
                 f"July views ({july_n} posts)",
                 size=10, color=PALETTE["muted"])

        add_rect(slide, x + 0.2, y + 2.0, card_w - 0.4, 0.015, PALETTE["rule"])
        add_text(slide, x + 0.2, y + 2.15, card_w - 0.4, 0.3,
                 "JUNE BASELINE", size=10, color=PALETTE["muted"], bold=True)
        add_text(slide, x + 0.2, y + 2.45, card_w - 0.4, 0.6,
                 fmt_int(june_v), size=24, color=PALETTE["ink"], bold=True)
        add_text(slide, x + 0.2, y + 3.1, card_w - 0.4, 0.3,
                 f"({june_n} posts)", size=10, color=PALETTE["muted"])

    slide_footer_takeaway(slide,
        "IG remains the workhorse. YouTube is the new signal — each episode is a compounding asset, not a spike.")


def slide_categories_july(prs, july_posts):
    """Category performance in July with expanded taxonomy."""
    slide = slide_new(prs)
    slide_header(slide, "Content taxonomy",
                 "What we posted in July, and how it performed",
                 "Average views per post across all platforms. Higher = more reach per piece.")

    buckets = group_by_category(july_posts)

    # Sort by total views desc, keep top 8 for a readable grid
    scored = []
    for cat, posts in buckets.items():
        if not posts:
            continue
        views = [int(p.get("view_count") or 0) for p in posts]
        scored.append({
            "cat": cat,
            "count": len(posts),
            "total": sum(views),
            "avg": sum(views) / len(views),
            "top": max(posts, key=lambda p: int(p.get("view_count") or 0)),
        })
    scored.sort(key=lambda s: -s["total"])
    top = scored[:8]

    # 4x2 grid
    card_w, card_h, gap = 2.9, 1.9, 0.15
    x0 = (LAYOUT_W - card_w * 4 - gap * 3) / 2
    y0 = 2.3

    label_map = {
        "trujillo_scandal":              ("Trujillo Scandal",       PALETTE["cherry"]),
        "lisette_scandal":               ("Lisette Scandal",        PALETTE["cherry"]),
        "ai_fake_news":                  ("AI Fake News",           PALETTE["cherry"]),
        "council_chaos":                 ("Council Chaos",          PALETTE["cherry"]),
        "ice_immigration":               ("ICE / Immigration",      PALETTE["cherry"]),
        "elections":                     ("Elections",              PALETTE["navy"]),
        "judicial":                      ("Judicial",               PALETTE["navy"]),
        "politics_other":                ("Politics (other)",       PALETTE["navy"]),
        "political_individual_highlight":("Political Highlight",    PALETTE["sky"]),
        "resident_highlight":            ("Resident Highlight",     PALETTE["green"]),
        "downtown_development":          ("Downtown Development",   PALETTE["gold"]),
        "american_pride":                ("American Pride",         PALETTE["cherry"]),
        "fifa_positive":                 ("FIFA — Positive",        PALETTE["green"]),
        "fifa_critique":                 ("FIFA — Critique",        PALETTE["muted"]),
        "food":                          ("Food",                   PALETTE["green"]),
        "brand_growth":                  ("Brand Growth",           PALETTE["blush"]),
        "community_events":              ("Community Events",       PALETTE["ink"]),
        "community_local":               ("Community Local",        PALETTE["ink"]),
        "other":                         ("Other",                  PALETTE["muted"]),
    }

    for i, s in enumerate(top):
        col = i % 4
        row = i // 4
        x = x0 + col * (card_w + gap)
        y = y0 + row * (card_h + gap)
        label, color = label_map.get(s["cat"], (s["cat"].title(), PALETTE["muted"]))

        add_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        add_rect(slide, x, y, 0.06, card_h, color)
        add_text(slide, x + 0.2, y + 0.15, card_w - 0.4, 0.3,
                 label.upper(), size=9, color=color, bold=True)
        add_text(slide, x + 0.2, y + 0.45, card_w - 0.4, 0.7,
                 f"{s['avg']:,.0f}", size=28, color=PALETTE["navy"],
                 bold=True, font="Georgia")
        add_text(slide, x + 0.2, y + 1.15, card_w - 0.4, 0.25,
                 "avg views/post", size=9, color=PALETTE["muted"])
        add_text(slide, x + 0.2, y + 1.42, card_w - 0.4, 0.5,
                 f"{s['count']} posts · {s['total']:,} total views",
                 size=9, color=PALETTE["ink"])

    slide_footer_takeaway(slide,
        "The taxonomy is broader in July. Each new bucket is a lane the sponsor can choose to invest more into.")


def slide_top_posts_july(prs, july_posts):
    """Top 8 individual posts in July by views, across all platforms."""
    slide = slide_new(prs)
    slide_header(slide, "Breakouts",
                 "Top July posts by views",
                 "Ranked across all platforms. These are the pieces that pulled outside the average.")

    posts = sorted(july_posts, key=lambda p: -int(p.get("view_count") or 0))[:8]

    col_widths = [0.5, 1.1, 6.4, 1.5, 1.3, 1.4]
    headers = ["#", "PLATFORM", "POST", "VIEWS", "LIKES", "CATEGORY"]
    row_y = 2.4
    row_h = 0.5
    x_base = 0.6

    add_rect(slide, x_base, row_y, sum(col_widths), row_h, PALETTE["navy"])
    cx = x_base
    for hd, w in zip(headers, col_widths):
        add_text(slide, cx + 0.1, row_y + 0.12, w - 0.2, 0.3,
                 hd, size=10, color=PALETTE["cream"], bold=True)
        cx += w

    for i, p in enumerate(posts, 1):
        ry = row_y + row_h + (i - 1) * row_h
        if i % 2 == 0:
            add_rect(slide, x_base, ry, sum(col_widths), row_h, PALETTE["soft"])
        if i <= 3:
            add_rect(slide, x_base, ry, 0.06, row_h, PALETTE["cherry"])

        title = (p.get("title") or p.get("description_excerpt") or "")[:75]
        plat = {"instagram": "IG", "tiktok": "TT", "facebook": "FB", "youtube": "YT"}.get(p.get("platform"), "?")
        cat = p.get("_category", "other").replace("_", " ").title()

        cx = x_base
        add_text(slide, cx + 0.1, ry + 0.12, col_widths[0] - 0.2, 0.3,
                 str(i), size=12, color=PALETTE["cherry"], bold=True)
        cx += col_widths[0]
        add_text(slide, cx + 0.1, ry + 0.12, col_widths[1] - 0.2, 0.3,
                 plat, size=11, color=PALETTE["navy"], bold=True)
        cx += col_widths[1]
        add_text(slide, cx + 0.1, ry + 0.12, col_widths[2] - 0.2, 0.3,
                 title, size=10, color=PALETTE["ink"])
        cx += col_widths[2]
        add_text(slide, cx + 0.1, ry + 0.12, col_widths[3] - 0.2, 0.3,
                 fmt_int(p.get("view_count")), size=12, color=PALETTE["navy"], bold=True)
        cx += col_widths[3]
        add_text(slide, cx + 0.1, ry + 0.12, col_widths[4] - 0.2, 0.3,
                 fmt_int(p.get("like_count")), size=11, color=PALETTE["ink"])
        cx += col_widths[4]
        add_text(slide, cx + 0.1, ry + 0.12, col_widths[5] - 0.2, 0.3,
                 cat[:22], size=9, color=PALETTE["muted"])

    slide_footer_takeaway(slide,
        "The mix at the top has expanded — Downey Growth and American Pride joined the political leaders.")


def slide_takeaway(prs, june_posts, july_posts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    add_rect(slide, 0, 4.7, LAYOUT_W, LAYOUT_H - 4.7, PALETTE["sky"])

    add_text(slide, 0.8, 0.6, 11, 0.4,
             "THE READ", size=12, color=PALETTE["sky"], bold=True)
    add_text(slide, 0.8, 1.0, 11.5, 1.0,
             "July diversified the show",
             size=36, color=PALETTE["white"], bold=True, font="Georgia")
    add_text(slide, 0.8, 2.15, 11.5, 0.6,
             "June was proof of concept — controversy engine at scale. "
             "July was the sequel — same engine, plus new lanes.",
             size=15, color=PALETTE["cream"])

    bullets = [
        f"June posts (Jun 8–30 window): {len(june_posts)}  →  July posts (Jul 1–31 window): {len(july_posts)}",
        "Political-only spine broadened to include Downey growth stories + American pride content",
        "Paul Granata episode (Jul 2) opened a new lane: local business + downtown revival",
        "July 4th content proved the audience responds to positive/celebratory as well",
        "Brand growth self-reporting posts landed — audience validated the meta-narrative",
    ]
    for i, b in enumerate(bullets):
        add_text(slide, 1.0, 3.2 + i * 0.32, 11.3, 0.3,
                 "—  " + b, size=12, color=PALETTE["cream"])

    add_text(slide, 0.8, 5.1, 11.5, 0.5,
             "WHAT TO ASK THE SPONSOR", size=12, color=PALETTE["ink"], bold=True)
    add_text(slide, 0.8, 5.5, 11.5, 1.0,
             "Where should we lean in August — heat, growth, or pride?",
             size=26, color=PALETTE["white"], bold=True, font="Georgia")


# ── main ────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path,
                    default=Path(__file__).parent / "MONTHLY_REPORT_2026-07.pptx")
    ap.add_argument("--kpis-root", type=Path,
                    default=Path(__file__).parent)
    args = ap.parse_args()

    snaps = load_all_snapshots(args.kpis_root / "episodes_public")
    if not snaps:
        sys.exit(f"No snapshots found under {args.kpis_root/'episodes_public'}")
    print(f"Loaded {len(snaps)} snapshots")

    june_posts = posts_in_window(snaps, JUNE_START, JUNE_END)
    july_posts = posts_in_window(snaps, JULY_START, JULY_END)
    print(f"  June window: {len(june_posts)} unique posts")
    print(f"  July window: {len(july_posts)} unique posts")

    prs = Presentation()
    prs.slide_width = Inches(LAYOUT_W)
    prs.slide_height = Inches(LAYOUT_H)

    slide_title(prs)
    slide_month_shift(prs, july_posts, june_posts)
    slide_platform_growth(prs, snaps)
    slide_categories_july(prs, july_posts)
    slide_top_posts_july(prs, july_posts)
    slide_takeaway(prs, june_posts, july_posts)

    prs.save(args.out)
    print(f"✓ Wrote: {args.out}")
    print(f"  {len(prs.slides)} slides")
    print(f"  Open with: open '{args.out}'")


if __name__ == "__main__":
    main()
