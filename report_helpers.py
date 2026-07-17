"""Talking Downey — shared building blocks for KPI report decks.

Extracted so build_july_report.py, build_overall_report.py, and
build_sponsor_content_report.py don't each re-implement the palette,
slide primitives, categorization, and snapshot loading logic. The
June builder (build_june_report.py) stays self-contained by design —
it was the frozen baseline before the July expansion.

Public surface:
  PALETTE, LAYOUT_W, LAYOUT_H — visual constants
  RGB, add_rect, add_text, fmt_int — pptx primitives
  normalize_date, snap_pt_label — date helpers
  categorize_v2 — expanded taxonomy with subcategories
  load_all_snapshots — scan every episode folder, return chronological snapshots
  posts_in_window — flatten records across snapshots, dedupe by URL, filter by date
"""
from __future__ import annotations

import json
import re as _re
from datetime import date, datetime, timedelta, timezone
from pathlib import Path

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


PT_OFFSET_HOURS = -7
LAYOUT_W, LAYOUT_H = 13.333, 7.5

# Engagement start — everything before this is prior editor's work, not ours.
ENGAGEMENT_START = "2026-06-08"


def RGB(r, g, b):
    return RGBColor(r, g, b)


PALETTE = {
    "cherry":  RGB(0x99, 0x00, 0x11),
    "navy":    RGB(0x2F, 0x3C, 0x7E),
    "ink":     RGB(0x14, 0x18, 0x33),
    "cream":   RGB(0xFC, 0xF6, 0xF5),
    "white":   RGB(0xFF, 0xFF, 0xFF),
    "muted":   RGB(0x6B, 0x70, 0x82),
    "rule":    RGB(0xE2, 0xDE, 0xDC),
    "soft":    RGB(0xF5, 0xEC, 0xEA),
    "green":   RGB(0x1B, 0x7A, 0x43),
    "gold":    RGB(0xC9, 0xA2, 0x27),   # positive/american
    "sky":     RGB(0x2E, 0x86, 0xAB),   # growth
    "blush":   RGB(0xE8, 0x7E, 0x8F),   # brand awareness
}


# ── slide primitives ────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size, color=None, bold=False,
             align="left", font="Helvetica Neue"):
    color = color or PALETTE["ink"]
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tb


def fmt_int(n):
    if n is None:
        return "—"
    try:
        return f"{int(n):,}"
    except (ValueError, TypeError):
        return "—"


def normalize_date(raw) -> str | None:
    if raw is None:
        return None
    s = str(raw)
    if not s:
        return None
    if len(s) == 8 and s.isdigit():
        return f"{s[0:4]}-{s[4:6]}-{s[6:8]}"
    if "-" in s and len(s) >= 10:
        return s[:10]
    try:
        t = int(s)
        if t > 10**9:
            return datetime.fromtimestamp(t, tz=timezone.utc).date().isoformat()
    except (ValueError, TypeError):
        pass
    return None


def snap_pt_label(snap) -> str:
    raw = (snap or {}).get("pulled_at", "")
    try:
        dt = datetime.fromisoformat(raw.replace("Z", "+00:00"))
        dt_pt = dt + timedelta(hours=PT_OFFSET_HOURS)
        return dt_pt.strftime("%b %-d (%a)")
    except (ValueError, TypeError, AttributeError):
        return (snap or {}).get("snapshot_tag", "?")


# ── expanded taxonomy ───────────────────────────────────────────────────
#
# Categories = broad topic buckets used across the whole report.
# Subcategories = sponsor-facing granular tags for the sponsor deck.
#
# Order matters: regex checked top-to-bottom, first hit wins. Put the
# most specific/high-signal patterns FIRST (trujillo scandal beats
# generic "political").
#
CATEGORY_ORDER = [
    # Scandal-first: critical/investigative content about specific figures.
    # Huntington Park lives here — Trujillo is HP mayor, HP posts are almost always about his conduct/policies.
    ("trujillo_scandal",              r"\b(trujillo|mario a\.?\s*trujillo|huntington park|huntington|hp mayor)\b"),
    ("lisette_scandal",               r"\b(lisette|lizsette|liz\s*sette)\b"),
    ("ai_fake_news",                  r"\b(ai\s*fake|fake\s*(page|news|logo)|slander|deepfake|ai\s*video)\b"),
    # Chaos = incident-driven only. Plain "council" mentions fall through to elections/political_individual_highlight.
    ("council_chaos",                 r"\b(council chambers|chambers.*(escort|yelling|police|outburst|warned)|f-bomb|f bomb|elderly resident.*f off|escorted out.*meeting|outburst.*council|forceful.*escort|yelling.*chambers|police escort.*council)\b"),
    ("ice_immigration",               r"\b(ice|federal investigation|immigrant|nicaragua)\b"),
    # Judicial BEFORE political_individual_highlight so judges stay judges.
    ("judicial",                      r"\b(judge|judicial|prosecutor|courtroom|jury|homicide|attorney|criminal|prop 50|pat connolly|maria ghobadi|gloria marin|irene lee)\b"),
    # Personality-profile buckets — NOT scandal, just highlights.
    ("political_individual_highlight", r"\b(carrie|uva|hector sosa|hector de la torre|erik lutz|mario a\.?\s*guerra|nader moghaddam|meet the|inside the campaign|reelect(ion)?|running for city council|for downey city council|district 4|willing to appoint)\b"),
    ("resident_highlight",            r"\b(granata|gately|30 years|restaurant.*owner|owner.*restaurant|deacon|father's story|diabetes|small business owner|one 12 caffe|smile on|running a restaurant)\b"),
    ("elections",                     r"\b(vote|voting|election|campaign|candidate|measure er|endorsement|primary)\b"),
    ("politics_other",                r"\b(mayor|politician|democrat|republican|committee|fined|congress|city leader)\b"),
    ("american_pride",                r"\b(american|independence|declaration|founding fathers|1,337 words|july 4|4th of july|memorial|veteran|patriot|embarrassed to be american|pride in this country)\b"),
    # FIFA BEFORE downtown_development so Stonewood watch parties land in FIFA.
    ("fifa_positive",                 r"\b(watch part|25,000|stonewood mall.*fifa|unite|passion|celebrate responsibly|championship game|sports unite)\b"),
    ("fifa_critique",                 r"\b(fifa.*money|billion|money-making|substitution|crowding out|profit.*fifa|world cup.*tax)\b"),
    # Physical growth of downtown — places, retail, infrastructure.
    ("downtown_development",          r"\b(round one|round1|stonewood|new asian|new restaurant|opening|main street|ymca|in-n-out|bowling|downtown improvement|gateway sign|columbia space|redevelop|new building|new business)\b"),
    ("food",                          r"\b(food|cuisine|restaurant|asian|dish|chef|eat|dining|bbq|pho|korean|mama lu|porto|fun box|chula|gyu|mango|brooklyn square|pizza|breakfast|lunch)\b"),
    ("brand_growth",                  r"\b(social media growth|exponential|100k views|54k|growth.*combined|views.*combined|our.*growth|social.*growth)\b"),
    ("community_events",              r"\b(pageant|night market|chamber|event|festival|holy week|lent|holiday|graffiti)\b"),
    ("community_local",               r"\b(downey|pico rivera|cudahy|city seal|20 years|letter)\b"),
]


def categorize_v2(text: str) -> str:
    t = (text or "").lower()
    for name, pat in CATEGORY_ORDER:
        if _re.search(pat, t):
            return name
    return "other"


# Display grouping for the sponsor deck — collapses fine-grained tags
# into presentable buckets. The sponsor sees "TRUTH & SCANDAL" not
# "trujillo_scandal · lisette_scandal · ai_fake_news · council_chaos".
SPONSOR_GROUPS = [
    ("Truth & Scandal Reporting",       ["trujillo_scandal", "lisette_scandal", "ai_fake_news", "council_chaos", "ice_immigration"],
     "cherry", "Investigative — the heat that drives peak reach."),
    ("Elections & Judicial",             ["elections", "judicial", "politics_other"],
     "navy",   "Coverage of campaigns, courts, and local government."),
    ("Resident Highlights",              ["resident_highlight"],
     "green",  "Non-political citizen profiles — small business owners, community members. Granata, Gately, the deacon."),
    ("Political Individual Highlights",  ["political_individual_highlight"],
     "sky",    "Candidate/elected personality profiles — Carrie Uva, Hector Sosa. Not scandal-driven; just meet-the-person."),
    ("Downtown Development",             ["downtown_development"],
     "gold",   "Physical growth of Downey — Round One, Stonewood, YMCA, new businesses opening."),
    ("American Pride",                   ["american_pride"],
     "cherry", "Patriotic content — July 4th, veterans, pride pieces."),
    ("Food & Local Flavor",              ["food"],
     "green",  "Restaurants, chefs, dining — the everyday Downey."),
    ("FIFA / World Cup",                 ["fifa_positive", "fifa_critique"],
     "muted",  "Sports commentary — split between celebration and critique."),
    ("Brand Growth Content",             ["brand_growth"],
     "blush",  "Self-referential updates — 'our growth is exponential' style posts."),
    ("Community & Culture",              ["community_events", "community_local"],
     "ink",    "Local color, events, city life."),
    ("Other / Uncategorized",            ["other"],
     "muted",  "Content that doesn't fit the above buckets — usually casual/BTS."),
]


# ── snapshot loading ────────────────────────────────────────────────────

def load_all_snapshots(episodes_root: Path) -> list[dict]:
    """Scan every episode_public/*/snapshots/*.json, return sorted by pulled_at."""
    snaps = []
    for ep_dir in sorted(episodes_root.iterdir()):
        if not ep_dir.is_dir():
            continue
        snap_dir = ep_dir / "snapshots"
        if not snap_dir.is_dir():
            continue
        for f in sorted(snap_dir.glob("*.json")):
            if f.name.startswith("."):
                continue
            try:
                snaps.append(json.loads(f.read_text()))
            except Exception:
                pass
    snaps.sort(key=lambda s: s.get("pulled_at", ""))
    return snaps


def posts_in_window(snaps: list[dict], start_iso: str, end_iso: str | None = None,
                    platforms: tuple = ("instagram", "tiktok", "facebook", "youtube")) -> list[dict]:
    """Return one record per unique URL, from posts uploaded in [start, end).

    Merges across every snapshot — later snapshots win (they have fresh view
    counts). Filters by upload_date within window. `end_iso` inclusive of
    that day if provided; None means open-ended.
    """
    latest_by_url: dict[str, dict] = {}
    for snap in snaps:
        for r in snap.get("records", []):
            url = r.get("url")
            if url:
                latest_by_url[url] = r
    out = []
    start_d = date.fromisoformat(start_iso)
    end_d = date.fromisoformat(end_iso) if end_iso else None
    for url, r in latest_by_url.items():
        plat = r.get("platform")
        if plat not in platforms:
            continue
        nd = normalize_date(r.get("upload_date"))
        if not nd:
            continue
        try:
            d = date.fromisoformat(nd)
        except ValueError:
            continue
        if d < start_d:
            continue
        if end_d and d > end_d:
            continue
        r_copy = dict(r)
        r_copy["_upload_date_iso"] = nd
        r_copy["_category"] = categorize_v2(
            (r.get("title") or "") + " " + (r.get("description_excerpt") or "")
        )
        out.append(r_copy)
    return out


def group_by_category(posts: list[dict]) -> dict[str, list[dict]]:
    """Bucket posts by their _category tag."""
    buckets: dict[str, list[dict]] = {}
    for p in posts:
        cat = p.get("_category", "other")
        buckets.setdefault(cat, []).append(p)
    return buckets


# ── standard slide chrome ───────────────────────────────────────────────

def slide_new(prs, bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, bg_color or PALETTE["cream"])
    return slide


def slide_header(slide, eyebrow: str, title: str, subtitle: str = ""):
    add_text(slide, 0.8, 0.5, 11.7, 0.4,
             eyebrow.upper(), size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             title, size=30, color=PALETTE["navy"], bold=True, font="Georgia")
    if subtitle:
        add_text(slide, 0.8, 1.75, 11.7, 0.4,
                 subtitle, size=12, color=PALETTE["muted"])


def slide_footer_takeaway(slide, text: str):
    add_text(slide, 0.8, 6.85, 11.7, 0.5,
             text, size=12, color=PALETTE["ink"], bold=True)
