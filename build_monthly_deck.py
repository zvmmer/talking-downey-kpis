"""Talking Downey — MONTHLY rollup deck.

Walks every episode under `episodes_public/` for a given month and builds
ONE consolidated PPTX showing the full body of work — not per-episode
silos.

Why this exists:
    Per-episode reports are useful for one client meeting. The client wants
    to see the WHOLE MONTH of output and how it performed cumulatively.
    This script gives that view.

Deck structure:
    1.  COVER         — "Talking Downey — June 2026" + headline stats
    2.  AT A GLANCE   — total posts, total views, top performer, ep count
    3.  PER EPISODE   — 1 section per episode (multiple slides if needed)
    4.  PLATFORM ROLL — top-N across all episodes per platform
    5.  WORK SHIPPED  — list of episodes published this month

Usage:
    MCP/.venv/bin/python Projects/talking_downey/kpis/build_monthly_deck.py \\
        --month 2026-06 \\
        --out Projects/talking_downey/kpis/MONTHLY_REPORT_2026-06.pptx
"""
from __future__ import annotations

import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

# Reuse all the slide-building helpers + palette from the per-episode deck.
sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_report_deck import (  # noqa: E402
    PALETTE, LAYOUT_W, LAYOUT_H,
    add_filled_rect, add_text, fmt_int, normalize_date,
    posts_on, platform_post_growth, sum_views, top_n_by_views,
    build_platform_full_table_slide,
)
from pptx import Presentation  # noqa: E402

EPISODES_DIR = Path(__file__).parent / "episodes_public"

MONTH_NAMES = {
    1: "January", 2: "February", 3: "March", 4: "April",
    5: "May", 6: "June", 7: "July", 8: "August",
    9: "September", 10: "October", 11: "November", 12: "December",
}

PLATFORM_LABELS = {
    "youtube": "YouTube",
    "tiktok": "TikTok",
    "instagram": "Instagram",
    "facebook": "Facebook",
}


# =============================================================================
# Episode discovery
# =============================================================================

def discover_episodes(month: str) -> list[dict]:
    """Return episode dicts for every <date>_<slug>/ folder posted in the given
    month (YYYY-MM). Each dict carries the latest snapshot's records."""
    out = []
    if not EPISODES_DIR.exists():
        return out
    for ep_dir in sorted(EPISODES_DIR.iterdir()):
        if not ep_dir.is_dir():
            continue
        name = ep_dir.name
        # Expect <YYYY-MM-DD>_<slug>/
        if not name[:7] == month:
            continue
        snaps_dir = ep_dir / "snapshots"
        if not snaps_dir.exists():
            continue
        snaps = sorted(snaps_dir.glob("*.json"))
        if not snaps:
            continue
        # Load all snapshots (chronological by mtime, since tag ordering is
        # alphabetical and t10 would sort before t2 — use mtime as tie-break).
        all_snaps = []
        for s in sorted(snaps, key=lambda p: p.stat().st_mtime):
            try:
                all_snaps.append(json.loads(s.read_text()))
            except Exception:
                continue
        if not all_snaps:
            continue
        latest = all_snaps[-1]
        out.append({
            "dir": ep_dir,
            "title": latest.get("episode", name),
            "posted_at": latest.get("posted_at", name[:10]),
            "snapshots": all_snaps,
            "latest": latest,
        })
    # Sort by posted_at descending — newest episode first
    out.sort(key=lambda e: e["posted_at"], reverse=True)
    return out


# =============================================================================
# Aggregations
# =============================================================================

def aggregate_month(episodes: list[dict]) -> dict:
    """Roll up across all episodes for headline stats."""
    total_posts = 0
    total_views = 0
    total_likes = 0
    total_comments = 0
    by_platform = {}
    all_records = []
    for ep in episodes:
        recs = ep["latest"].get("records", [])
        for r in recs:
            total_posts += 1
            total_views += int(r.get("view_count") or 0)
            total_likes += int(r.get("like_count") or 0)
            total_comments += int(r.get("comment_count") or 0)
            plat = r.get("platform", "?")
            p = by_platform.setdefault(plat, {"posts": 0, "views": 0, "likes": 0})
            p["posts"] += 1
            p["views"] += int(r.get("view_count") or 0)
            p["likes"] += int(r.get("like_count") or 0)
            all_records.append(r)
    return {
        "episode_count": len(episodes),
        "total_posts": total_posts,
        "total_views": total_views,
        "total_likes": total_likes,
        "total_comments": total_comments,
        "by_platform": by_platform,
        "all_records": all_records,
    }


def top_n_across_month(all_records: list[dict], n: int = 5) -> list[dict]:
    return sorted(
        [r for r in all_records if r.get("view_count")],
        key=lambda r: r.get("view_count") or 0,
        reverse=True,
    )[:n]


# =============================================================================
# Slides
# =============================================================================

def build_cover(prs, month: str, agg: dict, episodes: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["navy"])

    # Header strip
    yyyy, mm = month.split("-")
    month_label = f"{MONTH_NAMES[int(mm)]} {yyyy}"

    add_text(slide, 0.8, 0.7, 12, 0.4,
             "TALKING DOWNEY", size=14, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 1.15, 12, 1.0,
             f"Monthly Report · {month_label}",
             size=42, color=PALETTE["cream"], bold=True, font="Georgia")
    add_text(slide, 0.8, 2.3, 12, 0.5,
             f"Across {agg['episode_count']} episode{'s' if agg['episode_count'] != 1 else ''} · {agg['total_posts']} posts · {fmt_int(agg['total_views'])} views",
             size=18, color=PALETTE["cream"])

    # Big stat boxes
    box_y = 3.6
    box_h = 2.4
    box_w = 2.9
    gap = 0.2
    start_x = 0.8
    stats = [
        ("EPISODES", str(agg["episode_count"]), "this month"),
        ("POSTS", fmt_int(agg["total_posts"]), "across platforms"),
        ("VIEWS", fmt_int(agg["total_views"]), "cumulative"),
        ("LIKES", fmt_int(agg["total_likes"]), "engagement"),
    ]
    for i, (label, big, sub) in enumerate(stats):
        x = start_x + i * (box_w + gap)
        add_filled_rect(slide, x, box_y, box_w, box_h, PALETTE["cream"])
        add_text(slide, x + 0.2, box_y + 0.25, box_w - 0.4, 0.4,
                 label, size=11, color=PALETTE["cherry"], bold=True)
        add_text(slide, x + 0.2, box_y + 0.7, box_w - 0.4, 1.0,
                 big, size=36, color=PALETTE["navy"], bold=True, font="Georgia")
        add_text(slide, x + 0.2, box_y + 1.75, box_w - 0.4, 0.4,
                 sub, size=11, color=PALETTE["muted"])

    add_text(slide, 0.8, 6.3, 12, 0.4,
             f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} · auto-updates as new posts ship",
             size=10, color=PALETTE["cream"])


def build_top_performers(prs, month: str, agg: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 12, 0.4,
             "TOP PERFORMERS", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 12, 0.9,
             f"Best 5 posts this month, any platform",
             size=28, color=PALETTE["navy"], bold=True, font="Georgia")

    top = top_n_across_month(agg["all_records"], n=5)
    if not top:
        add_text(slide, 0.8, 3.0, 12, 0.5,
                 "No posts captured yet.", size=14, color=PALETTE["muted"])
        return

    row_y = 2.0
    row_h = 0.95
    for i, r in enumerate(top, 1):
        ry = row_y + (i - 1) * row_h
        # Row background — zebra
        if i % 2 == 0:
            add_filled_rect(slide, 0.8, ry, 11.7, row_h, PALETTE["white"])
        # Rank badge
        add_filled_rect(slide, 0.8, ry + 0.15, 0.6, row_h - 0.3, PALETTE["cherry"])
        add_text(slide, 0.8, ry + 0.18, 0.6, row_h - 0.3,
                 f"#{i}", size=18, color=PALETTE["cream"], bold=True, align="center")
        # Platform + title
        plat = (r.get("platform") or "?").upper()
        title = r.get("title", "")[:80]
        add_text(slide, 1.6, ry + 0.1, 7.5, 0.35,
                 plat, size=10, color=PALETTE["cherry"], bold=True)
        add_text(slide, 1.6, ry + 0.4, 7.5, 0.5,
                 title, size=14, color=PALETTE["navy"], bold=True)
        # Right-side stats
        add_text(slide, 9.3, ry + 0.1, 3.2, 0.4,
                 "VIEWS", size=9, color=PALETTE["cherry"], bold=True, align="right")
        add_text(slide, 9.3, ry + 0.35, 3.2, 0.6,
                 fmt_int(r.get("view_count") or 0),
                 size=22, color=PALETTE["navy"], bold=True, font="Georgia", align="right")


def build_episode_section_intro(prs, ep: dict, agg_for_ep: dict):
    """One slide introducing the episode + headline stats for it."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 12, 0.4,
             f"EPISODE · {ep['posted_at']}", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 12, 1.5,
             ep["title"], size=28, color=PALETTE["navy"], bold=True, font="Georgia")

    # Three stat tiles
    box_y = 2.8
    box_h = 2.0
    box_w = 3.8
    gap = 0.3
    start_x = 0.8
    stats = [
        ("POSTS", str(agg_for_ep["total_posts"]), "across platforms"),
        ("VIEWS", fmt_int(agg_for_ep["total_views"]), "cumulative"),
        ("ENGAGEMENT", fmt_int(agg_for_ep["total_likes"] + agg_for_ep["total_comments"]),
         "likes + comments"),
    ]
    for i, (label, big, sub) in enumerate(stats):
        x = start_x + i * (box_w + gap)
        add_filled_rect(slide, x, box_y, box_w, box_h, PALETTE["white"])
        add_text(slide, x + 0.25, box_y + 0.2, box_w - 0.5, 0.4,
                 label, size=11, color=PALETTE["cherry"], bold=True)
        add_text(slide, x + 0.25, box_y + 0.65, box_w - 0.5, 0.9,
                 big, size=32, color=PALETTE["navy"], bold=True, font="Georgia")
        add_text(slide, x + 0.25, box_y + 1.5, box_w - 0.5, 0.4,
                 sub, size=11, color=PALETTE["muted"])


def build_episode_posts_paginated(prs, ep: dict, posts_per_slide: int = 6):
    """One or more slides listing every post for the episode.

    The per-episode build_report_deck.build_platform_full_table_slide caps at
    8 rows. Here we paginate properly so nothing gets hidden, regardless of
    how many posts exist.
    """
    snaps = ep["snapshots"]
    target_date = ep["posted_at"]
    # Pull all post-level data per platform, but collapse to a single ordered
    # list across platforms so the slide shows the episode holistically.
    all_posts = []
    for plat in ["youtube", "tiktok", "instagram", "facebook"]:
        posts, _ = platform_post_growth(snaps, plat, target_date)
        for p in posts:
            p["_platform"] = plat
            all_posts.append(p)
    if not all_posts:
        return
    # Sort: day-1 first, then by current views desc
    all_posts.sort(key=lambda p: (
        not p.get("is_day1", False),
        -(p.get("now_views") or 0),
    ))

    pages = [all_posts[i:i + posts_per_slide]
             for i in range(0, len(all_posts), posts_per_slide)]
    for page_idx, page_posts in enumerate(pages, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])
        suffix = f" ({page_idx}/{len(pages)})" if len(pages) > 1 else ""
        add_text(slide, 0.8, 0.5, 12, 0.4,
                 f"POSTS · {ep['posted_at']}{suffix}",
                 size=12, color=PALETTE["cherry"], bold=True)
        add_text(slide, 0.8, 0.9, 11, 0.9,
                 ep["title"], size=22, color=PALETTE["navy"], bold=True, font="Georgia")

        # Table headers
        col_widths = [1.0, 1.2, 5.5, 1.8, 1.4, 1.0]
        headers = ["", "Platform", "Post", "Views (was → now)", "More views", "More likes"]
        row_y = 2.0
        row_h = 0.5
        add_filled_rect(slide, 0.8, row_y, sum(col_widths), row_h, PALETTE["navy"])
        cx = 0.8
        for hd, w in zip(headers, col_widths):
            add_text(slide, cx + 0.12, row_y + 0.12, w - 0.18, 0.3,
                     hd.upper(), size=9, color=PALETTE["cream"], bold=True)
            cx += w

        for i, p in enumerate(page_posts, 1):
            ry = row_y + row_h + (i - 1) * row_h
            is_day1 = p.get("is_day1", False)
            if is_day1:
                add_filled_rect(slide, 0.8, ry, sum(col_widths), row_h, PALETTE["soft"])
                add_filled_rect(slide, 0.8, ry, 0.06, row_h, PALETTE["cherry"])
            elif i % 2 == 0:
                add_filled_rect(slide, 0.8, ry, sum(col_widths), row_h, PALETTE["white"])
            cx = 0.8
            # tag
            badge = "DAY 1" if is_day1 else ("NEW" if p.get("is_new") else f"#{(page_idx - 1) * posts_per_slide + i}")
            badge_color = PALETTE["cherry"] if (is_day1 or p.get("is_new")) else PALETTE["muted"]
            add_filled_rect(slide, cx + 0.12, ry + 0.1, col_widths[0] - 0.25, row_h - 0.2, badge_color)
            add_text(slide, cx + 0.12, ry + 0.13, col_widths[0] - 0.25, row_h - 0.25,
                     badge, size=9, color=PALETTE["cream"], bold=True, align="center")
            cx += col_widths[0]
            # platform
            add_text(slide, cx + 0.12, ry + 0.13, col_widths[1] - 0.18, row_h - 0.2,
                     PLATFORM_LABELS.get(p.get("_platform", "?"), p.get("_platform", "?").upper()),
                     size=11, color=PALETTE["navy"], bold=True)
            cx += col_widths[1]
            # title
            title = (p.get("title") or "")[:62]
            add_text(slide, cx + 0.12, ry + 0.13, col_widths[2] - 0.18, row_h - 0.2,
                     title, size=11, color=PALETTE["ink"])
            cx += col_widths[2]
            # views was → now
            was = p.get("first_views") or 0
            now = p.get("now_views") or 0
            add_text(slide, cx + 0.12, ry + 0.13, col_widths[3] - 0.18, row_h - 0.2,
                     f"{fmt_int(was)}  →  {fmt_int(now)}", size=11, color=PALETTE["ink"])
            cx += col_widths[3]
            # Δ views
            dv = (p.get("now_views") or 0) - (p.get("first_views") or 0)
            add_text(slide, cx + 0.12, ry + 0.13, col_widths[4] - 0.18, row_h - 0.2,
                     f"+{fmt_int(dv)}" if dv > 0 else "—",
                     size=11, color=PALETTE["cherry"] if dv > 0 else PALETTE["muted"], bold=True)
            cx += col_widths[4]
            # Δ likes
            dl = (p.get("now_likes") or 0) - (p.get("first_likes") or 0)
            add_text(slide, cx + 0.12, ry + 0.13, col_widths[5] - 0.18, row_h - 0.2,
                     f"+{fmt_int(dl)}" if dl > 0 else "—",
                     size=11, color=PALETTE["cherry"] if dl > 0 else PALETTE["muted"], bold=True)


def build_work_shipped(prs, month: str, episodes: list[dict]):
    """Closer slide — visual list of episodes shipped this month."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["navy"])

    yyyy, mm = month.split("-")
    month_label = f"{MONTH_NAMES[int(mm)]} {yyyy}"
    add_text(slide, 0.8, 0.5, 12, 0.4,
             "WORK SHIPPED", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 12, 1.0,
             f"Episodes published this month",
             size=28, color=PALETTE["cream"], bold=True, font="Georgia")

    if not episodes:
        add_text(slide, 0.8, 3.0, 12, 0.5,
                 "No episodes published yet.", size=14, color=PALETTE["cream"])
        return

    # List each episode with date + title
    row_y = 2.2
    row_h = 0.7
    for i, ep in enumerate(episodes, 1):
        ry = row_y + (i - 1) * row_h
        add_text(slide, 0.8, ry + 0.05, 1.5, 0.5,
                 ep["posted_at"], size=14, color=PALETTE["cherry"], bold=True, font="Georgia")
        add_text(slide, 2.5, ry + 0.05, 10, 0.5,
                 ep["title"], size=14, color=PALETTE["cream"], bold=True)
        # Underline
        add_filled_rect(slide, 0.8, ry + 0.6, 11.7, 0.02, PALETTE["cherry"])


# =============================================================================
# Main
# =============================================================================

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--month", required=True, help="YYYY-MM, e.g. 2026-06")
    ap.add_argument("--out", type=Path, default=None,
                    help="Output PPTX path (default: kpis/MONTHLY_REPORT_<month>.pptx)")
    ap.add_argument("--posts-per-slide", type=int, default=6,
                    help="Pagination size for per-episode post lists (default 6)")
    args = ap.parse_args()

    episodes = discover_episodes(args.month)
    if not episodes:
        print(f"No episodes found for month {args.month}.", file=sys.stderr)
        sys.exit(1)

    print(f"Month: {args.month}")
    print(f"Episodes found: {len(episodes)}")
    for ep in episodes:
        print(f"  · {ep['posted_at']}  {ep['title']}  ({len(ep['snapshots'])} snapshots)")

    agg = aggregate_month(episodes)
    print(f"Totals: {agg['total_posts']} posts · {agg['total_views']:,} views · "
          f"{agg['total_likes']:,} likes")

    prs = Presentation()
    prs.slide_width = int(LAYOUT_W * 914400)
    prs.slide_height = int(LAYOUT_H * 914400)

    # 1. Cover
    build_cover(prs, args.month, agg, episodes)
    # 2. Top performers
    build_top_performers(prs, args.month, agg)
    # 3. Per-episode sections
    for ep in episodes:
        agg_for_ep = {
            "total_posts": len(ep["latest"].get("records", [])),
            "total_views": sum(int(r.get("view_count") or 0)
                              for r in ep["latest"].get("records", [])),
            "total_likes": sum(int(r.get("like_count") or 0)
                              for r in ep["latest"].get("records", [])),
            "total_comments": sum(int(r.get("comment_count") or 0)
                                 for r in ep["latest"].get("records", [])),
        }
        build_episode_section_intro(prs, ep, agg_for_ep)
        build_episode_posts_paginated(prs, ep, posts_per_slide=args.posts_per_slide)
    # 4. Work shipped (closer)
    build_work_shipped(prs, args.month, episodes)

    out_path = args.out or (Path(__file__).parent / f"MONTHLY_REPORT_{args.month}.pptx")
    prs.save(out_path)
    print(f"\n✓ Wrote: {out_path}")
    print(f"  {len(prs.slides)} slides")
    print(f"  Open with: open '{out_path}'")


if __name__ == "__main__":
    main()
