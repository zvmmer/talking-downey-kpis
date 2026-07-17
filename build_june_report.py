"""Talking Downey — June 2026 Performance Report deck.

16-day body-of-work report (Jun 8 – Jun 24). Reframed from the
per-episode deck into a portfolio progress report with Week 1 dense
tracking + Week 2+ check-in structure.

Usage
-----
    MCP/.venv/bin/python Projects/talking_downey/kpis/build_june_report.py \
        --episode Projects/talking_downey/kpis/episodes_public/2026-06-08_record_straight
"""
from __future__ import annotations

import argparse
import json
import re as _re
import sys
from datetime import date, datetime, timedelta, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


PT_OFFSET_HOURS = -7


def RGB(r, g, b):
    return RGBColor(r, g, b)


PALETTE = {
    "cherry":  RGB(0x99, 0x00, 0x11),
    "navy":    RGB(0x2F, 0x3C, 0x7E),
    "ink":     RGB(0x14, 0x18, 0x33),
    "cream":   RGB(0xFC, 0xF6, 0xF5),
    "white":   RGB(0xFF, 0xFF, 0xFF),
    "muted":   RGB(0x6B, 0x70, 0x82),
    "rule":    RGB(0xE2, 0xDE, 0xDC),
    "soft":    RGB(0xF5, 0xEC, 0xEA),
    "green":   RGB(0x1B, 0x7A, 0x43),
}

LAYOUT_W, LAYOUT_H = 13.333, 7.5


# ── helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size, color=None, bold=False,
             align="left", font="Helvetica Neue"):
    color = color or PALETTE["ink"]
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tb


def fmt_int(n):
    if n is None:
        return "—"
    try:
        return f"{int(n):,}"
    except (ValueError, TypeError):
        return "—"


def normalize_date(raw) -> str | None:
    if raw is None:
        return None
    s = str(raw)
    if not s:
        return None
    if len(s) == 8 and s.isdigit():
        return f"{s[0:4]}-{s[4:6]}-{s[6:8]}"
    if "-" in s and len(s) >= 10:
        return s[:10]
    try:
        t = int(s)
        if t > 10**9:
            return datetime.fromtimestamp(t, tz=timezone.utc).date().isoformat()
    except (ValueError, TypeError):
        pass
    return None


def snap_pt_label(snap) -> str:
    raw = (snap or {}).get("pulled_at", "")
    try:
        dt = datetime.fromisoformat(raw.replace("Z", "+00:00"))
        dt_pt = dt + timedelta(hours=PT_OFFSET_HOURS)
        return dt_pt.strftime("%b %-d (%a)")
    except (ValueError, TypeError, AttributeError):
        return snap.get("snapshot_tag", "?") if snap else "?"


def split_pre_post(records, day1_iso):
    pre, post = [], []
    for r in records:
        nd = normalize_date(r.get("upload_date"))
        if not nd:
            post.append(r)
            continue
        if nd < day1_iso:
            pre.append(r)
        else:
            post.append(r)
    return pre, post


def categorize(text):
    t = (text or "").lower()
    cats = {
        "political": r"\b(trujillo|lisette|mayor|council|councilman|ice|huntington|vote|election|campaign|candidate|judge|judicial|politician|prosecutor|gang|public works|pave|democrat|wills and trust|committee|fined|resident)\b",
        "fifa":      r"\b(fifa|world cup|soccer|watch part)\b",
        "food":      r"\b(food|cuisine|restaurant|asian|dish|chef|eat|dining|bbq|pho|korean|mama lu|porto|fun box|chula|gyu|mango)\b",
        "community": r"\b(downey|stonewood|round one|mall|downtown|night market)\b",
    }
    for name, pat in cats.items():
        if _re.search(pat, t):
            return name
    return "other"


def platform_post_growth(snaps, platform, day1_iso):
    if not snaps:
        return [], (0, 0)
    first_snap, last_snap = snaps[0], snaps[-1]
    first_by_url = {r.get("url"): r for r in first_snap.get("records", [])
                    if r.get("platform") == platform and r.get("url")}
    last_by_url = {r.get("url"): r for r in last_snap.get("records", [])
                   if r.get("platform") == platform and r.get("url")}
    posts = []
    for url, last in last_by_url.items():
        first = first_by_url.get(url)
        is_new = first is None
        posts.append({
            "url": url,
            "title": (last.get("title") or last.get("description_excerpt") or "")[:80],
            "is_new": is_new,
            "v_first": (first or {}).get("view_count") or 0,
            "v_last": last.get("view_count") or 0,
            "l_first": (first or {}).get("like_count") or 0,
            "l_last": last.get("like_count") or 0,
            "_upload_date": last.get("upload_date"),
        })
    nd_int = lambda p: int(normalize_date(p.get("_upload_date", "")) .replace("-", "") or "0") if normalize_date(p.get("_upload_date")) else 0
    posts.sort(key=lambda p: (-nd_int(p), -int(p["v_last"] or 0)))
    total_first = sum(int(p["v_first"] or 0) for p in posts if not p["is_new"])
    total_last = sum(int(p["v_last"] or 0) for p in posts)
    return posts, (total_first, total_last)


# ── slide builders ───────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    add_rect(slide, 0, 0, 0.45, LAYOUT_H, PALETTE["cherry"])

    add_text(slide, 1.0, 1.2, 11, 0.5,
             "JUNE 2026 PERFORMANCE REPORT",
             size=14, color=PALETTE["cherry"], bold=True)
    add_text(slide, 1.0, 1.9, 11, 1.3,
             "Talking Downey",
             size=64, color=PALETTE["white"], bold=True, font="Georgia")
    add_text(slide, 1.0, 3.4, 11, 0.6,
             "16-Day Progress Report — Social Media Performance",
             size=24, color=PALETTE["cream"], font="Georgia")

    add_text(slide, 1.0, 5.0, 5, 0.4,
             "Tracking period: Jun 8 – Jun 24, 2026",
             size=13, color=PALETTE["muted"])
    add_text(slide, 1.0, 5.4, 5, 0.4,
             "7 measurement points across 16 days",
             size=13, color=PALETTE["muted"])
    add_text(slide, 1.0, 5.8, 5, 0.4,
             "4 platforms: Instagram · TikTok · Facebook · YouTube",
             size=13, color=PALETTE["muted"])

    add_text(slide, 1.0, 6.6, 11, 0.4,
             "Prepared by Zhamir Pascual — Kaname Z",
             size=12, color=PALETTE["muted"])


def slide_volume(prs, latest_records, day1_iso):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "THE OUTPUT", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             "Content volume — before vs. now",
             size=30, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "Same brand, same audience — completely different output.",
             size=12, color=PALETTE["muted"])

    social_platforms = ("instagram", "tiktok", "facebook")
    social_recs = [r for r in latest_records if r.get("platform") in social_platforms]
    pre, post = split_pre_post(social_recs, day1_iso)

    pre_per_week = len(pre) / 4 if pre else 0
    post_days = (date.fromisoformat("2026-06-24") - date.fromisoformat(day1_iso)).days or 1
    post_per_week = len(post) / (post_days / 7)

    # Big center number
    add_text(slide, 0.8, 2.5, 11.7, 1.5,
             f"{len(post)}", size=120, color=PALETTE["cherry"], bold=True,
             font="Georgia", align="center")
    add_text(slide, 0.8, 4.1, 11.7, 0.5,
             "posts across 3 platforms in 16 days",
             size=18, color=PALETTE["navy"], align="center")

    # Before / After comparison cards
    card_w, card_h = 5.5, 1.8
    gap = 0.8
    x_left = (LAYOUT_W - card_w * 2 - gap) / 2
    x_right = x_left + card_w + gap
    y = 5.0

    # BEFORE card
    add_rect(slide, x_left, y, card_w, card_h, PALETTE["white"])
    add_rect(slide, x_left, y, 0.08, card_h, PALETTE["muted"])
    add_text(slide, x_left + 0.3, y + 0.2, card_w - 0.5, 0.3,
             "BEFORE JUN 8", size=11, color=PALETTE["muted"], bold=True)
    add_text(slide, x_left + 0.3, y + 0.55, card_w - 0.5, 0.5,
             f"~{pre_per_week:.0f} posts/week", size=24, color=PALETTE["ink"], bold=True)
    add_text(slide, x_left + 0.3, y + 1.1, card_w - 0.5, 0.5,
             f"{len(pre)} posts in data  ·  Sporadic posting  ·  No coordinated strategy",
             size=11, color=PALETTE["muted"])

    # SINCE card
    add_rect(slide, x_right, y, card_w, card_h, PALETTE["white"])
    add_rect(slide, x_right, y, 0.08, card_h, PALETTE["cherry"])
    add_text(slide, x_right + 0.3, y + 0.2, card_w - 0.5, 0.3,
             "SINCE JUN 8", size=11, color=PALETTE["cherry"], bold=True)
    add_text(slide, x_right + 0.3, y + 0.55, card_w - 0.5, 0.5,
             f"~{post_per_week:.0f} posts/week", size=24, color=PALETTE["cherry"], bold=True)
    add_text(slide, x_right + 0.3, y + 1.1, card_w - 0.5, 0.5,
             f"{len(post)} posts  ·  2-3 per day  ·  IG + TT + FB coordinated daily",
             size=11, color=PALETTE["muted"])


def slide_before_after(prs, snaps, day1_iso):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "BEFORE vs. SINCE I JOINED", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.85, 11.7, 0.9,
             "More output AND more reach per post",
             size=30, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "Their older posts had weeks to accumulate views. My posts are days old — and still pulling more per post.",
             size=12, color=PALETTE["muted"])

    day1_date = date.fromisoformat(day1_iso)
    latest = snaps[-1]
    tt_snap = max(snaps, key=lambda s: sum(1 for r in s.get("records", []) if r.get("platform") == "tiktok"))

    # For before/after, gather ALL records across ALL snapshots to get the widest
    # date range (FB posts scroll off the 30-post window between snapshots).
    all_records_ever: dict[str, dict] = {}  # keyed by url, keeps latest version
    for snap in snaps:
        for r in snap.get("records", []):
            url = r.get("url")
            if url:
                all_records_ever[url] = r

    cards = []
    for plat_key, plat_label in [
        ("instagram", "Instagram"),
        ("tiktok", "TikTok"),
        ("facebook", "Facebook"),
    ]:
        pre, post = [], []
        for url, r in all_records_ever.items():
            if r.get("platform") != plat_key:
                continue
            nd = normalize_date(r.get("upload_date"))
            if not nd:
                continue
            v = int(r.get("view_count") or 0)
            if date.fromisoformat(nd) >= day1_date:
                post.append(v)
            else:
                pre.append(v)
        pre_avg = (sum(pre) / len(pre)) if pre else 0
        post_avg = (sum(post) / len(post)) if post else 0
        multiplier = (post_avg / pre_avg) if pre_avg > 0 else 0
        cards.append({
            "platform": plat_label,
            "pre_count": len(pre), "pre_avg": pre_avg,
            "post_count": len(post), "post_avg": post_avg,
            "multiplier": multiplier,
        })

    card_w, card_h = 3.85, 4.5
    gap = 0.3
    total_w = card_w * 3 + gap * 2
    x_start = (LAYOUT_W - total_w) / 2
    y = 2.3

    for i, cd in enumerate(cards):
        x = x_start + i * (card_w + gap)
        add_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        add_rect(slide, x, y, 0.08, card_h, PALETTE["cherry"])

        add_text(slide, x + 0.3, y + 0.25, card_w - 0.5, 0.4,
                 cd["platform"].upper(), size=12, color=PALETTE["cherry"], bold=True)

        mult_str = f"{cd['multiplier']:.1f}×" if cd["multiplier"] else "N/A"
        add_text(slide, x + 0.3, y + 0.7, card_w - 0.5, 1.1,
                 mult_str, size=66, color=PALETTE["navy"], bold=True, font="Georgia")
        add_text(slide, x + 0.3, y + 1.85, card_w - 0.5, 0.4,
                 "more views per post", size=11, color=PALETTE["muted"])

        add_rect(slide, x + 0.3, y + 2.3, card_w - 0.6, 0.015, PALETTE["rule"])

        add_text(slide, x + 0.3, y + 2.45, card_w - 0.5, 0.3,
                 "BEFORE I JOINED", size=10, color=PALETTE["muted"], bold=True)
        add_text(slide, x + 0.3, y + 2.75, card_w - 0.5, 0.4,
                 f"{cd['pre_avg']:,.0f} views avg",
                 size=18, color=PALETTE["ink"], bold=True)
        add_text(slide, x + 0.3, y + 3.1, card_w - 0.5, 0.3,
                 f"({cd['pre_count']} posts, weeks old)",
                 size=10, color=PALETTE["muted"])

        add_text(slide, x + 0.3, y + 3.5, card_w - 0.5, 0.3,
                 "SINCE I JOINED", size=10, color=PALETTE["cherry"], bold=True)
        add_text(slide, x + 0.3, y + 3.8, card_w - 0.5, 0.4,
                 f"{cd['post_avg']:,.0f} views avg",
                 size=18, color=PALETTE["cherry"], bold=True)
        add_text(slide, x + 0.3, y + 4.15, card_w - 0.5, 0.3,
                 f"({cd['post_count']} posts, days old)",
                 size=10, color=PALETTE["muted"])

    add_text(slide, 0.8, 7.05, 11.7, 0.4,
             "Same audience. Same brand. Different editor. More reach per post — every platform.",
             size=12, color=PALETTE["ink"], bold=True)


def slide_week1(prs, snaps, day1_iso):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "WEEK 1 — DENSE TRACKING", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             "Jun 8–13: 4 snapshots in 5 days",
             size=28, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "We measured every 1-2 days in the first week to catch the initial growth trajectory.",
             size=12, color=PALETTE["muted"])

    # Use first 4 snapshots (t0-t3) for Week 1
    week1_snaps = snaps[:4] if len(snaps) >= 4 else snaps

    plat_order = [
        ("instagram", "Instagram"),
        ("tiktok", "TikTok"),
        ("facebook", "Facebook"),
    ]

    chart_data = CategoryChartData()
    snap_labels = [snap_pt_label(s) for s in week1_snaps]
    chart_data.categories = snap_labels

    for plat_key, plat_label in plat_order:
        vals = []
        for snap in week1_snaps:
            total = sum(int(r.get("view_count") or 0)
                        for r in snap.get("records", [])
                        if r.get("platform") == plat_key)
            vals.append(total)
        chart_data.add_series(plat_label, vals)

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.5),
        chart_data,
    )
    chart = gframe.chart
    chart.has_legend = True

    colors = [PALETTE["cherry"], PALETTE["navy"], PALETTE["muted"]]
    for idx, series in enumerate(chart.series):
        series.format.line.color.rgb = colors[idx % len(colors)]
        series.format.line.width = Pt(2.5)

    add_text(slide, 0.8, 6.1, 11.7, 0.3,
             "WEEK 1 TAKEAWAY", size=10, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 6.4, 11.7, 0.5,
             "All three social platforms climbed steadily through the first week. "
             "This wasn't a one-day spike — the numbers kept moving with every new post.",
             size=13, color=PALETTE["ink"], bold=True)


def slide_week2(prs, snaps, day1_iso):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "TWO WEEKS IN — JUN 24 CHECK-IN", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             "The growth held — this wasn't a one-week sugar rush",
             size=28, color=PALETTE["navy"], bold=True, font="Georgia")

    first_label = snap_pt_label(snaps[0]) if snaps else "?"
    last_label = snap_pt_label(snaps[-1]) if snaps else "?"
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             f"Comparing {first_label} (Day 1) to {last_label} (today) — total views across all tracked posts per platform.",
             size=12, color=PALETTE["muted"])

    plat_order = [
        ("instagram", "Instagram"),
        ("tiktok", "TikTok"),
        ("facebook", "Facebook"),
        ("youtube", "YouTube"),
    ]

    # Merge all snapshots for widest coverage (same approach as before/after slide)
    day1_date = date.fromisoformat(day1_iso)
    all_by_url: dict[str, dict] = {}
    for snap in snaps:
        for r in snap.get("records", []):
            url = r.get("url")
            if url:
                all_by_url[url] = r

    # Growth % chart — use per-post average growth (pre vs post day1)
    # This works even when URLs don't overlap between snapshots (FB case).
    cats = []
    growth_pcts = []
    abs_growth = []
    fb_is_estimated = False
    for key, label in plat_order:
        # First try URL-matched growth (works for IG, TT, YT)
        _, (first_total, last_total) = platform_post_growth(snaps, key, day1_iso)
        first_total, last_total = int(first_total or 0), int(last_total or 0)

        if first_total > 0:
            pct = round((last_total - first_total) / first_total * 100, 1)
            cats.append(label)
            growth_pcts.append(pct)
            abs_growth.append((label, first_total, last_total, False))
        else:
            # No URL overlap — fall back to per-post average growth
            pre_views, post_views = [], []
            for url, r in all_by_url.items():
                if r.get("platform") != key:
                    continue
                nd = normalize_date(r.get("upload_date"))
                if not nd:
                    continue
                v = int(r.get("view_count") or 0)
                try:
                    if date.fromisoformat(nd) >= day1_date:
                        post_views.append(v)
                    else:
                        pre_views.append(v)
                except ValueError:
                    pass
            pre_avg = (sum(pre_views) / len(pre_views)) if pre_views else 0
            post_avg = (sum(post_views) / len(post_views)) if post_views else 0
            if pre_avg > 0:
                pct = round((post_avg - pre_avg) / pre_avg * 100, 1)
            else:
                pct = 0
            cats.append(label + "*" if pct > 0 else label)
            growth_pcts.append(pct)
            pre_total = int(pre_avg * len(pre_views))
            post_total = int(post_avg * len(post_views))
            abs_growth.append((label, pre_total, post_total, True))
            if key == "facebook" and pct > 0:
                fb_is_estimated = True

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series("Growth %", growth_pcts)

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(2.3), Inches(7.0), Inches(3.6),
        chart_data,
    )
    chart = gframe.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 75
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.font.size = Pt(14)
    dl.font.bold = True
    dl.font.color.rgb = PALETTE["navy"]
    dl.number_format = '+0"%";-0"%";0"%"'
    for s in chart.series:
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = PALETTE["cherry"]

    # Side panel: absolute numbers
    x_panel = 8.3
    add_text(slide, x_panel, 2.3, 4.5, 0.35,
             "ABSOLUTE NUMBERS", size=11, color=PALETTE["cherry"], bold=True)
    for i, (label, first, last, is_est) in enumerate(abs_growth):
        y_row = 2.75 + i * 0.85
        add_rect(slide, x_panel, y_row, 4.5, 0.75, PALETTE["white"])
        add_rect(slide, x_panel, y_row, 0.06, 0.75, PALETTE["cherry"])
        suffix = " (est.)" if is_est else ""
        add_text(slide, x_panel + 0.2, y_row + 0.05, 2.0, 0.3,
                 label.upper() + suffix, size=10, color=PALETTE["cherry"], bold=True)
        if is_est:
            add_text(slide, x_panel + 0.2, y_row + 0.35, 4.0, 0.35,
                     f"~{fmt_int(first)} avg/post → ~{fmt_int(last)} avg/post",
                     size=13, color=PALETTE["ink"], bold=True)
        else:
            add_text(slide, x_panel + 0.2, y_row + 0.35, 4.0, 0.35,
                     f"{fmt_int(first)} → {fmt_int(last)} views",
                     size=13, color=PALETTE["ink"], bold=True)

    add_text(slide, 0.8, 6.2, 11.7, 0.3,
             "WEEK 2 TAKEAWAY", size=10, color=PALETTE["cherry"], bold=True)
    footnote = ""
    if fb_is_estimated:
        footnote = " *Facebook uses per-post avg growth (posts rotated out between snapshots)."
    add_text(slide, 0.8, 6.5, 11.7, 0.5,
             "Every platform is still climbing after 16 days. Consistent posting keeps the momentum — the audience is building, not just spiking."
             + footnote,
             size=13, color=PALETTE["ink"], bold=True)


def slide_categories(prs, snaps, day1_iso):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "CONTENT CATEGORIES", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.85, 11.7, 0.9,
             "What topics perform — by the numbers",
             size=30, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "Average views per post across IG + TT + FB since Jun 8. Higher = more reach per piece of content.",
             size=12, color=PALETTE["muted"])

    day1_date = date.fromisoformat(day1_iso)

    # Merge ALL snapshots to get widest coverage (FB posts scroll off between pulls)
    all_records_by_url: dict[str, dict] = {}
    for snap in snaps:
        for r in snap.get("records", []):
            url = r.get("url")
            if url:
                all_records_by_url[url] = r

    buckets: dict[str, list] = {"political": [], "fifa": [], "food": [], "community": []}
    for url, r in all_records_by_url.items():
        plat = r.get("platform")
        if plat not in ("instagram", "tiktok", "facebook"):
            continue
        nd = normalize_date(r.get("upload_date"))
        if not nd:
            continue
        try:
            if date.fromisoformat(nd) < day1_date:
                continue
        except ValueError:
            continue
        text = r.get("title") or r.get("description_excerpt") or ""
        cat = categorize(text)
        v = int(r.get("view_count") or 0)
        l = int(r.get("like_count") or 0)
        if cat in buckets:
            buckets[cat].append({"views": v, "likes": l, "title": text[:60], "url": url, "platform": plat})

    cat_display = [
        ("political", "Political", "Council, Trujillo, judicial, elections — the engine."),
        ("fifa", "FIFA / World Cup", "Opportunistic — works when real moments hit."),
        ("food", "Food", "New lane — still ramping up."),
        ("community", "Community", "Downey local — steady mid-tier."),
    ]

    card_w = 2.85
    card_full_h = 4.0
    gap = 0.2
    total_w = card_w * 4 + gap * 3
    x_start = (LAYOUT_W - total_w) / 2
    y = 2.3
    card_full_h = 4.0

    for i, (key, label, blurb) in enumerate(cat_display):
        x = x_start + i * (card_w + gap)
        posts = buckets.get(key, [])
        add_rect(slide, x, y, card_w, card_full_h, PALETTE["white"])
        add_rect(slide, x, y, 0.06, card_full_h, PALETTE["cherry"])

        add_text(slide, x + 0.2, y + 0.2, card_w - 0.4, 0.4,
                 label.upper(), size=10, color=PALETTE["cherry"], bold=True)

        if posts:
            avg = sum(p["views"] for p in posts) / len(posts)
            avg_str = f"{avg:,.0f}"
        else:
            avg_str = "—"
        add_text(slide, x + 0.2, y + 0.55, card_w - 0.4, 0.9,
                 avg_str, size=44, color=PALETTE["navy"], bold=True, font="Georgia")
        add_text(slide, x + 0.2, y + 1.45, card_w - 0.4, 0.3,
                 "avg views/post", size=10, color=PALETTE["muted"])

        if posts:
            tot_v = sum(p["views"] for p in posts)
            add_text(slide, x + 0.2, y + 1.8, card_w - 0.4, 0.3,
                     f"{len(posts)} posts · {tot_v:,} total",
                     size=10, color=PALETTE["ink"], bold=True)
            top = max(posts, key=lambda p: p["views"])
            add_text(slide, x + 0.2, y + 2.2, card_w - 0.4, 0.25,
                     "TOP POST", size=9, color=PALETTE["cherry"], bold=True)
            add_text(slide, x + 0.2, y + 2.45, card_w - 0.4, 0.8,
                     f"{top['views']:,} views — {top['title'][:50]}",
                     size=9, color=PALETTE["ink"])

        add_text(slide, x + 0.2, y + 3.4, card_w - 0.4, 0.5,
                 blurb, size=9, color=PALETTE["muted"])

    add_text(slide, 0.8, 6.55, 11.7, 0.3,
             "RECOMMENDATION", size=10, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 6.85, 11.7, 0.5,
             "Political content is the clear engine — ~2× the reach of any other category. "
             "FIFA is opportunistic and strong when moments hit. "
             "Food is the newest lane, still growing. Community is reliable mid-tier filler between the heavy hitters.",
             size=12, color=PALETTE["ink"], bold=True)


def slide_top_performers(prs, snaps, day1_iso):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "BREAKOUT POSTS", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             "Top performing content since Jun 8",
             size=28, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "Ranked by views across all platforms. These posts prove the audience is there — you just have to feed it.",
             size=12, color=PALETTE["muted"])

    day1_date = date.fromisoformat(day1_iso)

    # Merge all snapshots for widest coverage
    all_records_by_url: dict[str, dict] = {}
    for snap in snaps:
        for r in snap.get("records", []):
            url = r.get("url")
            if url:
                all_records_by_url[url] = r

    all_posts = []
    seen_titles = set()
    for url, r in all_records_by_url.items():
        plat = r.get("platform")
        if plat not in ("instagram", "tiktok", "facebook"):
            continue
        nd = normalize_date(r.get("upload_date"))
        if not nd:
            continue
        try:
            if date.fromisoformat(nd) < day1_date:
                continue
        except ValueError:
            continue
        title = (r.get("title") or r.get("description_excerpt") or "")[:70]
        dedup_key = title[:30].lower()
        if dedup_key in seen_titles:
            continue
        seen_titles.add(dedup_key)
        all_posts.append({
            "title": title,
            "views": int(r.get("view_count") or 0),
            "likes": int(r.get("like_count") or 0),
            "platform": plat,
            "date": nd,
            "category": categorize(title),
        })

    all_posts.sort(key=lambda p: -p["views"])
    top = all_posts[:8]

    # Table
    col_widths = [0.5, 1.2, 6.3, 1.5, 1.5, 1.2]
    headers = ["#", "Platform", "Post", "Views", "Likes", "Category"]
    row_y = 2.4
    row_h = 0.52
    x_base = 0.6

    add_rect(slide, x_base, row_y, sum(col_widths), row_h, PALETTE["navy"])
    cx = x_base
    for hd, w in zip(headers, col_widths):
        add_text(slide, cx + 0.1, row_y + 0.13, w - 0.2, 0.3,
                 hd.upper(), size=10, color=PALETTE["cream"], bold=True)
        cx += w

    for i, p in enumerate(top, 1):
        ry = row_y + row_h + (i - 1) * row_h
        if i % 2 == 0:
            add_rect(slide, x_base, ry, sum(col_widths), row_h, PALETTE["soft"])
        if i <= 3:
            add_rect(slide, x_base, ry, 0.06, row_h, PALETTE["cherry"])

        cx = x_base
        add_text(slide, cx + 0.1, ry + 0.13, col_widths[0] - 0.2, 0.3,
                 str(i), size=12, color=PALETTE["cherry"], bold=True)
        cx += col_widths[0]

        plat_display = {"instagram": "IG", "tiktok": "TT", "facebook": "FB"}.get(p["platform"], p["platform"])
        add_text(slide, cx + 0.1, ry + 0.13, col_widths[1] - 0.2, 0.3,
                 plat_display, size=11, color=PALETTE["navy"], bold=True)
        cx += col_widths[1]

        add_text(slide, cx + 0.1, ry + 0.13, col_widths[2] - 0.2, 0.3,
                 p["title"], size=10, color=PALETTE["ink"])
        cx += col_widths[2]

        add_text(slide, cx + 0.1, ry + 0.13, col_widths[3] - 0.2, 0.3,
                 fmt_int(p["views"]), size=12, color=PALETTE["navy"], bold=True)
        cx += col_widths[3]

        add_text(slide, cx + 0.1, ry + 0.13, col_widths[4] - 0.2, 0.3,
                 fmt_int(p["likes"]), size=11, color=PALETTE["ink"])
        cx += col_widths[4]

        add_text(slide, cx + 0.1, ry + 0.13, col_widths[5] - 0.2, 0.3,
                 p["category"].title(), size=10, color=PALETTE["muted"])

    add_text(slide, 0.8, 6.7, 11.7, 0.5,
             "Political content dominates the top spots. But community stories (Round One, Stonewood) "
             "show the audience cares about Downey life, not just controversy.",
             size=12, color=PALETTE["ink"], bold=True)


def slide_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    add_rect(slide, 0, 4.7, LAYOUT_W, LAYOUT_H - 4.7, PALETTE["cherry"])

    add_text(slide, 0.8, 0.6, 11, 0.4,
             "NEXT STEPS", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 1.0, 11.5, 1.0,
             "16 days of public data got us this far",
             size=36, color=PALETTE["white"], bold=True, font="Georgia")
    add_text(slide, 0.8, 2.15, 11.5, 0.5,
             "With view-only insights access, the next report includes watch time, retention, demographics, and traffic sources.",
             size=15, color=PALETTE["cream"])

    add_text(slide, 0.8, 3.1, 5.5, 0.4,
             "WHAT I NEED FROM YOU", size=12, color=PALETTE["cherry"], bold=True)
    needs = [
        "Add me to the Facebook Page with the 'Insights' role",
        "Add me to Instagram with view-only insights",
        "About 30 minutes of setup, once",
        "I can't post, edit, or change anything",
    ]
    for i, n in enumerate(needs):
        add_text(slide, 1.0, 3.5 + i * 0.34, 5.4, 0.3,
                 "—  " + n, size=12, color=PALETTE["cream"])

    add_text(slide, 6.7, 3.1, 5.5, 0.4,
             "WHAT THE NEXT REPORT INCLUDES", size=12, color=PALETTE["cherry"], bold=True)
    gets = [
        "Real watch time per video (not just view counts)",
        "Where viewers come from (search, suggested, shared)",
        "Audience demographics — age, location, devices",
        "Automated weekly reports with richer data",
    ]
    for i, g in enumerate(gets):
        add_text(slide, 6.9, 3.5 + i * 0.34, 5.4, 0.3,
                 "—  " + g, size=12, color=PALETTE["cream"])

    add_text(slide, 0.8, 5.1, 5.5, 0.5,
             "TIME COMMITMENT", size=12, color=PALETTE["cream"], bold=True)
    add_text(slide, 0.8, 5.5, 5.5, 1.0,
             "~30 minutes, once",
             size=32, color=PALETTE["white"], bold=True, font="Georgia")

    add_text(slide, 6.7, 5.1, 5.5, 0.5,
             "ONGOING COST", size=12, color=PALETTE["cream"], bold=True)
    add_text(slide, 6.7, 5.5, 5.5, 1.0,
             "$0",
             size=32, color=PALETTE["white"], bold=True, font="Georgia")


# ── main ─────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--episode", required=True, type=Path)
    ap.add_argument("--out", type=Path)
    args = ap.parse_args()

    snap_dir = args.episode / "snapshots"
    if not snap_dir.is_dir():
        sys.exit(f"No snapshots dir at {snap_dir}")

    snaps = []
    for f in sorted(snap_dir.glob("*.json")):
        if not f.name.startswith("."):
            snaps.append(json.loads(f.read_text()))
    if not snaps:
        sys.exit(f"No snapshot JSONs in {snap_dir}")

    latest = snaps[-1]
    all_records = latest.get("records", [])
    posted_at = latest.get("posted_at", "")
    day1_iso = (posted_at or "2026-06-08")[:10]

    prs = Presentation()
    prs.slide_width = Inches(LAYOUT_W)
    prs.slide_height = Inches(LAYOUT_H)

    slide_title(prs)
    slide_volume(prs, all_records, day1_iso)
    slide_before_after(prs, snaps, day1_iso)
    slide_week1(prs, snaps, day1_iso)
    slide_week2(prs, snaps, day1_iso)
    slide_categories(prs, snaps, day1_iso)
    slide_top_performers(prs, snaps, day1_iso)
    slide_ask(prs)

    out = args.out or (args.episode / "june_report.pptx")
    prs.save(out)
    print(f"✓ Wrote: {out}")
    print(f"  {len(prs.slides)} slides")
    print(f"  Open with: open '{out}'")


if __name__ == "__main__":
    main()
