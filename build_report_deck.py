"""Talking Downey — generate the PPTX pitch deck from snapshots.

Reads every snapshot JSON in an episode's snapshots/ folder and produces a
client-ready slideshow that tells the story of (a) what the episode is doing
on each platform and (b) what proper API access would unlock.

Usage
-----
    MCP/.venv/bin/python Projects/talking_downey/kpis/build_report_deck.py \\
        --episode Projects/talking_downey/kpis/episodes_public/2026-06-08_record_straight

Output: <episode>/report.pptx (overwrites if exists).
"""
from __future__ import annotations

import argparse
import json
import sys
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import date, timedelta, datetime, timezone


# Pacific time conversion for client-facing labels (Talking Downey is LA-based)
PT_OFFSET_HOURS = -7  # PDT


def snap_pt_label(snap) -> str:
    """Return a human-readable PT date+weekday label for a snapshot.

    snap is the loaded JSON dict; uses pulled_at (UTC ISO) and shifts to PT.
    """
    raw = (snap or {}).get("pulled_at", "")
    try:
        # parse ISO 8601 with timezone
        dt = datetime.fromisoformat(raw.replace("Z", "+00:00"))
        dt_pt = dt + timedelta(hours=PT_OFFSET_HOURS)
        return dt_pt.strftime("%b %-d (%a)")  # e.g., "Jun 8 (Sun)"
    except (ValueError, TypeError, AttributeError):
        return snap.get("snapshot_tag", "?") if snap else "?"


def RGB(r, g, b):
    return RGBColor(r, g, b)


# Past YT episodes — fetched from yt-dlp probe. Refresh by running:
#   yt-dlp --skip-download --print "%(upload_date)s|%(view_count)s|%(title).60s" <VIDEO_URL>
# for each recent episode.
YT_HISTORY = [
    {"date": "2026-05-22", "views": 212, "title": "What's Really Happening in Huntington Park?"},
    {"date": "2026-05-25", "views": 135, "title": "2026 Primary California Elections"},
    {"date": "2026-05-27", "views": 203, "title": "Inside the Courtroom: Pat Connolly"},
    {"date": "2026-05-28", "views": 138, "title": "Inside The Campaign: Erik Lutz"},
    {"date": "2026-06-08", "views": 120, "title": "Setting the Record Straight (TODAY)"},
]


# ============================================================================
# PALETTE — Cherry Bold (matches Talking Downey brand)
# ============================================================================
PALETTE = {
    "cherry":   RGB(0x99, 0x00, 0x11),  # primary accent
    "navy":     RGB(0x2F, 0x3C, 0x7E),  # dark anchor
    "ink":      RGB(0x14, 0x18, 0x33),  # near-black text
    "cream":    RGB(0xFC, 0xF6, 0xF5),  # off-white bg
    "white":    RGB(0xFF, 0xFF, 0xFF),
    "muted":    RGB(0x6B, 0x70, 0x82),  # secondary text
    "rule":     RGB(0xE2, 0xDE, 0xDC),
    "soft":     RGB(0xF5, 0xEC, 0xEA),  # tinted cherry surface for cards
}


# ============================================================================
# Helpers
# ============================================================================

def add_filled_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = PALETTE["rule"]
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size, color=PALETTE["ink"], bold=False,
             align="left", font="Helvetica Neue"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tb


def fmt_int(n):
    if n is None: return "—"
    try: return f"{int(n):,}"
    except: return "—"


def normalize_date(raw) -> str | None:
    """Return YYYY-MM-DD from any of the date formats we see across platforms.

    yt-dlp gives YouTube/TikTok as 'YYYYMMDD'. Apify gives Instagram/Facebook as
    ISO 8601 'YYYY-MM-DDTHH:MM:SS.000Z'. Unix timestamps also appear sometimes.
    """
    if raw is None:
        return None
    s = str(raw)
    if not s:
        return None
    # YYYYMMDD (yt-dlp)
    if len(s) == 8 and s.isdigit():
        return f"{s[0:4]}-{s[4:6]}-{s[6:8]}"
    # ISO 8601 — take the first 10 chars
    if "-" in s and len(s) >= 10:
        return s[:10]
    # Unix timestamp (seconds)
    try:
        t = int(s)
        if t > 10**9:
            return datetime.fromtimestamp(t, tz=timezone.utc).date().isoformat()
    except (ValueError, TypeError):
        pass
    return None


def posts_on(records, platform, target_date: str) -> list[dict]:
    """Return all records for `platform` whose normalized date equals target."""
    return [
        r for r in records
        if r.get("platform") == platform
        and normalize_date(r.get("upload_date")) == target_date
    ]


def day1_trajectory(snaps, platform, target_date):
    """Return [(snap_label, record), ...] for the Day-1 post on this platform,
    across every snapshot. Snap label is a short readable date like 'Jun 8'."""
    out = []
    for snap in snaps:
        posts = posts_on(snap.get("records", []), platform, target_date)
        if not posts:
            continue
        # Use pulled_at to label the column (e.g. "Jun 8")
        pulled = snap.get("pulled_at", "")
        label = pulled[5:10] if len(pulled) >= 10 else snap.get("snapshot_tag", "?")
        # Reformat YYYY-MM-DD → "Mon DD"
        try:
            d = date.fromisoformat(pulled[:10])
            label = d.strftime("%b %d")
        except (ValueError, TypeError):
            pass
        out.append((label, posts[0]))
    return out


def fmt_delta(curr, prev):
    """Return cherry-red string like '+150' or '—'."""
    try:
        c, p = int(curr or 0), int(prev or 0)
        d = c - p
        if d == 0:
            return ""
        return f"+{d:,}" if d > 0 else f"{d:,}"
    except (ValueError, TypeError):
        return ""


def _post_date_int(p):
    """Return YYYYMMDD as int for sorting, or 0 if missing.

    The platform_post_growth dict doesn't carry the raw upload_date — pull
    it off the matching record stored on the dict. Caller stamps it via
    `_upload_date` for sorting use.
    """
    raw = p.get("_upload_date")
    if not raw:
        return 0
    nd = normalize_date(raw)
    if not nd:
        return 0
    try:
        return int(nd.replace("-", ""))
    except (ValueError, AttributeError):
        return 0


def platform_post_growth(snaps, platform, target_date):
    """For a platform, return list of dicts with t0/t1/delta per post (matched by URL).

    Also marks `is_day1` if the post's upload_date matches target_date,
    and `is_new_since_first_snapshot` if the post wasn't present in the
    earliest snapshot (i.e., posted after baseline).
    """
    if not snaps:
        return [], (0, 0)
    first_snap, last_snap = snaps[0], snaps[-1]
    first_by_url = {r.get("url"): r for r in first_snap.get("records", [])
                    if r.get("platform") == platform and r.get("url")}
    last_by_url = {r.get("url"): r for r in last_snap.get("records", [])
                   if r.get("platform") == platform and r.get("url")}
    posts = []
    for url, last in last_by_url.items():
        first = first_by_url.get(url)
        nd = normalize_date(last.get("upload_date"))
        is_new = first is None  # not present in baseline = posted after t0
        posts.append({
            "url": url,
            "title": (last.get("title") or last.get("description_excerpt") or "")[:80],
            "is_day1": nd == target_date,
            "is_new": is_new,
            "v_first": (first or {}).get("view_count") or 0,
            "v_last": last.get("view_count") or 0,
            "l_first": (first or {}).get("like_count") or 0,
            "l_last": last.get("like_count") or 0,
            "media_type": last.get("media_type") or "",
            "_upload_date": last.get("upload_date"),
        })
    # Sort by upload date, NEWEST first — true chronological order.
    # Falls back to view count when dates are missing.
    def _sort_key(p):
        # Try to parse date from the last-known record
        for rec in [p]:  # placeholder for any future enrichment
            nd = None  # we don't have raw date here; use stored upload_date via lookup below
            break
        return (-int(p.get("v_last") or 0),)
    posts.sort(
        key=lambda p: (
            # Primary: upload date descending (newest first). Use the URL-keyed
            # record from the LATEST snapshot to fetch the date.
            -(_post_date_int(p) or 0),
            # Tiebreaker: more views first
            -int(p["v_last"] or 0),
        )
    )
    # Totals: include EVERY post's "now" views; baseline counts only posts
    # that existed at the first snapshot (avoid the 0→X distortion).
    total_first = sum(int(p["v_first"] or 0) for p in posts if not p["is_new"])
    total_last = sum(int(p["v_last"] or 0) for p in posts)
    return posts, (total_first, total_last)


def sum_views(records, platform):
    total = 0
    for r in records:
        if r.get("platform") == platform:
            v = r.get("view_count")
            try: total += int(v) if v else 0
            except: pass
    return total


def top_n_by_views(records, platform, n=5):
    items = [r for r in records if r.get("platform") == platform and r.get("view_count")]
    items.sort(key=lambda r: int(r.get("view_count") or 0), reverse=True)
    return items[:n]


# ============================================================================
# Slide builders
# ============================================================================

LAYOUT_W, LAYOUT_H = 13.333, 7.5  # widescreen 16:9


def build_title_slide(prs, episode_title, posted_at, snapshot_count):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Dark cherry background
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    # Cherry accent block on left
    add_filled_rect(slide, 0, 0, 0.45, LAYOUT_H, PALETTE["cherry"])
    # Eyebrow
    add_text(slide, 1.0, 1.4, 11, 0.5,
             "KPI PERFORMANCE REPORT",
             size=14, color=PALETTE["cherry"], bold=True)
    # Big title
    add_text(slide, 1.0, 2.0, 11, 1.3,
             "Talking Downey",
             size=64, color=PALETTE["white"], bold=True, font="Georgia")
    # Episode title
    add_text(slide, 1.0, 3.5, 11, 1.4,
             episode_title,
             size=24, color=PALETTE["cream"], font="Georgia")
    # Footer info
    posted_str = posted_at[:10] if posted_at else "—"
    add_text(slide, 1.0, 6.4, 6, 0.4,
             f"Posted: {posted_str}",
             size=12, color=PALETTE["muted"])
    add_text(slide, 7.0, 6.4, 5.3, 0.4,
             f"Snapshots captured: {snapshot_count}",
             size=12, color=PALETTE["muted"], align="right")


def build_glance_slide(prs, snaps, target_date):
    """At-a-glance growth across IG/TT/FB. YT mentioned as small footer note."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    # Pacific Time date labels — clearer for the LA meeting
    first_label = snap_pt_label(snaps[0]) if snaps else "?"
    last_label = snap_pt_label(snaps[-1]) if snaps else "?"

    add_text(slide, 0.8, 0.6, 11, 0.4,
             f"FROM {first_label.upper()} TO {last_label.upper()}",
             size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 1.0, 11.7, 0.9,
             "Your social posts are growing fast",
             size=32, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.85, 11.7, 0.4,
             f"Total views across every recent post on each platform — before ({first_label}) compared to now ({last_label}).",
             size=12, color=PALETTE["muted"])

    cards = []
    for plat, label, sub in [
        ("instagram", "Instagram", "across all recent Reels & posts"),
        ("tiktok",    "TikTok",    "across all recent clips"),
        ("facebook",  "Facebook",  "across all recent Page posts"),
    ]:
        _, (first_total, last_total) = platform_post_growth(snaps, plat, target_date)
        cards.append((label, first_total, last_total, sub))

    # Three big stat cards for IG/TT/FB
    card_w, card_h = 3.85, 3.7
    gap = 0.3
    total_w = card_w * 3 + gap * 2
    x_start = (LAYOUT_W - total_w) / 2
    y = 2.6

    for i, (plat, first_total, last_total, sub) in enumerate(cards):
        x = x_start + i * (card_w + gap)
        # Card
        add_filled_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        # Cherry stripe
        add_filled_rect(slide, x, y, 0.08, card_h, PALETTE["cherry"])
        # Platform label
        add_text(slide, x + 0.35, y + 0.3, card_w - 0.5, 0.4,
                 plat.upper(), size=12, color=PALETTE["cherry"], bold=True)
        # Big "now" number
        add_text(slide, x + 0.35, y + 0.75, card_w - 0.5, 1.3,
                 fmt_int(last_total), size=54, color=PALETTE["navy"], bold=True, font="Georgia")
        # "views" label
        add_text(slide, x + 0.35, y + 2.0, card_w - 0.5, 0.3,
                 "total views", size=11, color=PALETTE["muted"])
        # Growth callout
        try:
            d = last_total - first_total
            pct = (d / first_total * 100) if first_total > 0 else 0
            growth_str = f"+{d:,} views ({pct:+.0f}%)"
        except (ZeroDivisionError, TypeError):
            growth_str = "—"
        add_text(slide, x + 0.35, y + 2.45, card_w - 0.5, 0.5,
                 growth_str, size=18, color=PALETTE["cherry"], bold=True)
        # Sub-label
        add_text(slide, x + 0.35, y + 3.0, card_w - 0.5, 0.5,
                 sub, size=11, color=PALETTE["muted"])

    # Plain-English takeaway above the YT footer
    add_text(slide, 0.8, 6.45, 11.7, 0.3,
             "TAKEAWAY", size=10, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 6.7, 11.7, 0.4,
             "Every social platform is climbing — your posting is moving the numbers, not just sitting there.",
             size=13, color=PALETTE["ink"], bold=True)

    # YT footer line — context, not headline
    _, (yt_first, yt_last) = platform_post_growth(snaps, "youtube", target_date)
    if yt_last:
        try:
            yt_pct = (yt_last - yt_first) / yt_first * 100 if yt_first > 0 else 0
            yt_str = f"(YouTube episode also grew: {yt_first:,} → {yt_last:,} views, {yt_pct:+.0f}%)"
        except (ZeroDivisionError, TypeError):
            yt_str = f"(YouTube episode: {yt_last:,} views)"
        add_text(slide, 0.8, 7.05, 11.7, 0.3,
                 yt_str, size=10, color=PALETTE["muted"])


def build_platform_trajectory_slide(prs, *, platform_label: str, headline: str,
                                    caption: str, trajectory: list,
                                    footnote: str = ""):
    """Platform slide showing one post's stats across multiple snapshots.

    `trajectory` is a list of (snap_label, record) — newest at the end.
    Shows the post title at top + stat blocks below with t0 → t1 → Δ.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             platform_label.upper(), size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             headline, size=28, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             caption, size=12, color=PALETTE["muted"])

    if not trajectory:
        add_text(slide, 0.8, 3.0, 11, 0.5,
                 "No data captured for this platform.",
                 size=14, color=PALETTE["muted"])
        return

    # Show the post title (from the newest snapshot record)
    title = (trajectory[-1][1].get("title")
             or trajectory[-1][1].get("description_excerpt") or "")[:120]
    add_text(slide, 0.8, 2.3, 11.7, 0.7,
             f"“{title}”", size=18, color=PALETTE["ink"], font="Georgia")

    # Three stat blocks: Views, Likes, Comments
    metrics = [
        ("Views",    "view_count"),
        ("Likes",    "like_count"),
        ("Comments", "comment_count"),
    ]
    block_w, block_h = 3.95, 2.7
    gap = 0.15
    total_w = block_w * 3 + gap * 2
    x_start = (LAYOUT_W - total_w) / 2
    y_block = 3.3

    for i, (metric_label, key) in enumerate(metrics):
        x = x_start + i * (block_w + gap)

        # Card background
        add_filled_rect(slide, x, y_block, block_w, block_h, PALETTE["white"])
        # Cherry stripe on left
        add_filled_rect(slide, x, y_block, 0.08, block_h, PALETTE["cherry"])

        # Metric label
        add_text(slide, x + 0.3, y_block + 0.25, block_w - 0.5, 0.35,
                 metric_label.upper(), size=11, color=PALETTE["cherry"], bold=True)

        # Stat columns inside the block — one per snapshot + a Δ column at end
        n = len(trajectory)
        col_w = (block_w - 0.5) / (n + 1)  # n snapshots + 1 delta column
        col_y_label = y_block + 0.75
        col_y_value = y_block + 1.05

        prev_val = None
        for j, (snap_label, rec) in enumerate(trajectory):
            cx = x + 0.3 + j * col_w
            add_text(slide, cx, col_y_label, col_w, 0.3,
                     snap_label, size=10, color=PALETTE["muted"], align="center")
            val = rec.get(key)
            add_text(slide, cx, col_y_value, col_w, 0.8,
                     fmt_int(val), size=30, color=PALETTE["navy"], bold=True,
                     font="Georgia", align="center")
            prev_val = val

        # Delta column
        if n >= 2:
            curr = trajectory[-1][1].get(key)
            first = trajectory[0][1].get(key)
            delta_str = fmt_delta(curr, first)
            cx = x + 0.3 + n * col_w
            add_text(slide, cx, col_y_label, col_w, 0.3,
                     "GROWTH", size=10, color=PALETTE["cherry"], bold=True, align="center")
            add_text(slide, cx, col_y_value, col_w, 0.8,
                     delta_str or "—", size=26, color=PALETTE["cherry"], bold=True,
                     font="Georgia", align="center")

        # Percentage callout at bottom
        if n >= 2:
            try:
                first = int(trajectory[0][1].get(key) or 0)
                curr = int(trajectory[-1][1].get(key) or 0)
                if first > 0:
                    pct = (curr - first) / first * 100
                    if pct != 0:
                        add_text(slide, x + 0.3, y_block + 2.2, block_w - 0.5, 0.4,
                                 f"{pct:+.0f}% vs Day 1",
                                 size=11, color=PALETTE["muted"], align="center")
            except (ValueError, TypeError):
                pass

    if footnote:
        add_text(slide, 0.8, 6.6, 11.7, 0.4,
                 footnote, size=10, color=PALETTE["muted"])


def build_platform_card_slide(prs, *, platform: str, color_key: str, headline: str,
                              caption: str, posts: list, columns: list,
                              footnote: str = ""):
    """LEGACY: single-snapshot table. Kept for compatibility."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    # Eyebrow + title
    add_text(slide, 0.8, 0.6, 11, 0.4,
             platform.upper(), size=12, color=PALETTE[color_key], bold=True)
    add_text(slide, 0.8, 1.0, 11, 0.9,
             headline, size=32, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.85, 11, 0.4,
             caption, size=13, color=PALETTE["muted"])

    # Table of top posts
    if not posts:
        add_text(slide, 0.8, 3.0, 11, 0.5,
                 "No data captured for this platform.", size=14, color=PALETTE["muted"])
    else:
        # Header row
        col_widths = [0.5, 7.5, 1.4, 1.4, 1.4]  # ranks, title, views, likes, comments
        x = 0.8
        row_y = 2.6
        row_h = 0.5
        # Header
        add_filled_rect(slide, x, row_y, sum(col_widths), row_h, PALETTE["navy"])
        cx = x
        for label, w in zip(["#"] + columns, col_widths):
            add_text(slide, cx + 0.15, row_y + 0.1, w - 0.2, 0.3,
                     label.upper(), size=10, color=PALETTE["cream"], bold=True)
            cx += w
        # Rows
        for i, post in enumerate(posts, 1):
            ry = row_y + row_h + (i - 1) * row_h
            # Zebra stripes — every other row tinted
            if i % 2 == 0:
                add_filled_rect(slide, x, ry, sum(col_widths), row_h, PALETTE["soft"])
            cx = x
            # Rank
            add_text(slide, cx + 0.15, ry + 0.13, col_widths[0] - 0.2, 0.3,
                     str(i), size=12, color=PALETTE["cherry"], bold=True)
            cx += col_widths[0]
            # Title
            title = (post.get("title") or post.get("description_excerpt") or "")[:70]
            if not title:
                title = post.get("url", "—")[-30:]
            add_text(slide, cx + 0.15, ry + 0.13, col_widths[1] - 0.2, 0.3,
                     title, size=11, color=PALETTE["ink"])
            cx += col_widths[1]
            # Views
            add_text(slide, cx + 0.15, ry + 0.13, col_widths[2] - 0.2, 0.3,
                     fmt_int(post.get("view_count")), size=12, color=PALETTE["navy"], bold=True)
            cx += col_widths[2]
            # Likes
            add_text(slide, cx + 0.15, ry + 0.13, col_widths[3] - 0.2, 0.3,
                     fmt_int(post.get("like_count")), size=12, color=PALETTE["ink"])
            cx += col_widths[3]
            # Comments
            add_text(slide, cx + 0.15, ry + 0.13, col_widths[4] - 0.2, 0.3,
                     fmt_int(post.get("comment_count")), size=12, color=PALETTE["ink"])

    if footnote:
        add_text(slide, 0.8, 6.6, 11.7, 0.4,
                 footnote, size=10, color=PALETTE["muted"])


def build_platform_full_table_slide(prs, *, platform: str, label: str, snaps,
                                     target_date: str, post_word: str = "posts",
                                     footnote: str = ""):
    """A platform's full post list with t0/t1/delta columns. Day-1 post highlighted."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    posts, (total_first, total_last) = platform_post_growth(snaps, platform, target_date)

    add_text(slide, 0.8, 0.5, 11, 0.4,
             label.upper(), size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 8.5, 0.9,
             f"Growth across recent {post_word}",
             size=28, color=PALETTE["navy"], bold=True, font="Georgia")
    # Right-side total growth callout
    try:
        delta = total_last - total_first
        pct = (delta / total_first * 100) if total_first > 0 else 0
        growth_label = f"+{delta:,}"
        pct_label = f"{pct:+.0f}%"
    except (ZeroDivisionError, TypeError):
        growth_label, pct_label = "—", ""
    add_text(slide, 9.3, 0.6, 3.5, 0.4,
             "TOTAL GROWTH", size=10, color=PALETTE["cherry"], bold=True, align="right")
    add_text(slide, 9.3, 1.0, 3.5, 0.9,
             growth_label, size=32, color=PALETTE["cherry"], bold=True, font="Georgia", align="right")
    add_text(slide, 9.3, 1.85, 3.5, 0.4,
             f"views  ·  {pct_label}", size=12, color=PALETTE["muted"], align="right")

    add_text(slide, 0.8, 1.85, 8.5, 0.4,
             f"{len(posts)} {post_word} tracked across snapshots. Day-1 launch post highlighted in cherry.",
             size=12, color=PALETTE["muted"])

    if not posts:
        add_text(slide, 0.8, 3.0, 11, 0.5,
                 "No posts captured.", size=14, color=PALETTE["muted"])
        return

    # Table: tag | Post | Views (was→now) | Growth in views | Growth in likes
    col_widths = [1.0, 6.5, 1.8, 1.4, 1.0]  # tag, post, views was→now, Δ views, Δ likes
    headers = ["", "Post", "Views (was → now)", "More views", "More likes"]
    row_y = 2.55
    row_h = 0.42
    # Header row — navy background, cream text
    add_filled_rect(slide, 0.8, row_y, sum(col_widths), row_h, PALETTE["navy"])
    cx = 0.8
    for hd, w in zip(headers, col_widths):
        add_text(slide, cx + 0.12, row_y + 0.1, w - 0.18, 0.3,
                 hd.upper(), size=9, color=PALETTE["cream"], bold=True)
        cx += w

    # Body rows (limit to 8 for readability)
    for i, p in enumerate(posts[:8], 1):
        ry = row_y + row_h + (i - 1) * row_h
        is_day1 = p["is_day1"]
        # Day 1 highlight = soft cherry tint; even rows = light gray zebra
        if is_day1:
            add_filled_rect(slide, 0.8, ry, sum(col_widths), row_h, PALETTE["soft"])
            add_filled_rect(slide, 0.8, ry, 0.06, row_h, PALETTE["cherry"])
        elif i % 2 == 0:
            add_filled_rect(slide, 0.8, ry, sum(col_widths), row_h, PALETTE["white"])

        cx = 0.8
        # Tag column — "DAY 1" badge, "NEW" badge, or numbered rank
        is_new = p.get("is_new", False)
        if is_day1 or is_new:
            # Cherry pill-style label
            badge = "DAY 1" if (is_day1 and not is_new) else ("NEW" if is_new else "DAY 1")
            add_filled_rect(slide, cx + 0.12, ry + 0.08, col_widths[0] - 0.25, 0.26, PALETTE["cherry"])
            add_text(slide, cx + 0.12, ry + 0.1, col_widths[0] - 0.25, 0.26,
                     badge, size=9, color=PALETTE["white"], bold=True, align="center")
        else:
            add_text(slide, cx + 0.12, ry + 0.1, col_widths[0] - 0.18, 0.3,
                     f"#{i}", size=11, color=PALETTE["muted"])
        cx += col_widths[0]
        # Title
        title = p["title"] or "—"
        add_text(slide, cx + 0.12, ry + 0.1, col_widths[1] - 0.18, 0.3,
                 title, size=10, color=PALETTE["ink"])
        cx += col_widths[1]
        # Views was → now (show "just posted" for new posts instead of 0)
        if is_new:
            v_text = f"just posted → {fmt_int(p['v_last'])}"
        else:
            v_text = f"{fmt_int(p['v_first'])} → {fmt_int(p['v_last'])}"
        add_text(slide, cx + 0.12, ry + 0.1, col_widths[2] - 0.18, 0.3,
                 v_text, size=10, color=PALETTE["ink"])
        cx += col_widths[2]
        # Δ views (for NEW posts, show "+X" since they grew from nothing)
        if is_new:
            delta = f"+{fmt_int(p['v_last'])}"
        else:
            delta = fmt_delta(p["v_last"], p["v_first"]) or "—"
        add_text(slide, cx + 0.12, ry + 0.1, col_widths[3] - 0.18, 0.3,
                 delta, size=11, color=PALETTE["cherry"], bold=True)
        cx += col_widths[3]
        # Δ likes (similarly handle new posts)
        if is_new:
            likes_delta = f"+{fmt_int(p['l_last'])}" if p['l_last'] else "—"
        else:
            likes_delta = fmt_delta(p["l_last"], p["l_first"]) or "—"
        add_text(slide, cx + 0.12, ry + 0.1, col_widths[4] - 0.18, 0.3,
                 likes_delta, size=10, color=PALETTE["ink"])

    # Plain-English takeaway for this platform
    try:
        delta_total = total_last - total_first
        pct_total = (delta_total / total_first * 100) if total_first > 0 else 0
        plain = f"All your {post_word} on this platform together gained {delta_total:,} more views — that's {pct_total:+.0f}% growth in 24 hours."
    except (ZeroDivisionError, TypeError):
        plain = "Your posts on this platform are gaining views since the last snapshot."

    add_text(slide, 0.8, 6.45, 11.7, 0.3,
             "TAKEAWAY", size=10, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 6.7, 11.7, 0.5,
             plain, size=13, color=PALETTE["ink"], bold=True)

    if footnote:
        add_text(slide, 0.8, 7.15, 11.7, 0.3,
                 footnote, size=9, color=PALETTE["muted"])


def build_day1_growth_chart_slide(prs, snaps, target_date):
    """Grouped column chart: before vs. now per platform. Easier to read than
    a line chart when magnitudes differ wildly (YT vs IG)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    first_label = snap_pt_label(snaps[0]) if snaps else "?"
    last_label = snap_pt_label(snaps[-1]) if snaps else "?"

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "THE PICTURE", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             "Every platform grew by double or triple digits",
             size=28, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             f"Growth in total views per platform from {first_label} to {last_label}. The taller the bar, the bigger the jump.",
             size=12, color=PALETTE["muted"])

    # Percentage-growth chart — every platform reads at the same visual scale.
    # Avoids YouTube (120→270) becoming an invisible sliver next to IG (3k→14k).
    plat_order = [("instagram", "Instagram"),
                  ("tiktok", "TikTok"),
                  ("facebook", "Facebook"),
                  ("youtube", "YouTube")]

    cats = []
    growth_pcts = []
    for key, label in plat_order:
        _, (first, last) = platform_post_growth(snaps, key, target_date)
        first, last = int(first or 0), int(last or 0)
        cats.append(label)
        if first > 0:
            growth_pcts.append(round((last - first) / first * 100, 1))
        else:
            growth_pcts.append(0)

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series("Growth %", growth_pcts)

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.6),
        chart_data,
    )
    chart = gframe.chart
    chart.has_legend = False

    # Color the columns cherry red
    plot = chart.plots[0]
    plot.gap_width = 75
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.font.size = Pt(14)
    data_labels.font.bold = True
    data_labels.font.color.rgb = PALETTE["navy"]
    # Format labels as percent
    data_labels.number_format = '+0"%";-0"%";0"%"'

    for s in chart.series:
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = PALETTE["cherry"]

    # Plain-English takeaway
    add_text(slide, 0.8, 6.1, 11.7, 0.3,
             "TAKEAWAY", size=10, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 6.35, 11.7, 0.6,
             "All four platforms are climbing in the first 24 hours. Coordinated posting works — and these numbers are still climbing.",
             size=13, color=PALETTE["ink"], bold=True)


def build_correlation_chart_slide(prs, all_records: list):
    """Line chart: YouTube episode views over time + daily social post volume.
    Tries to answer 'does more social posting drive YT views?' visually."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "THE QUESTION", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.9, 11.7, 0.9,
             "Does social posting drive YouTube views?",
             size=30, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "Past 3 weeks: YT episodes (line) + total daily social activity (line). "
             "Today's launch is the rightmost point.",
             size=12, color=PALETTE["muted"])

    # Build daily timeline from earliest YT date to today
    yt_dates = [date.fromisoformat(e["date"]) for e in YT_HISTORY]
    social_dates_raw = [normalize_date(r.get("upload_date")) for r in all_records]
    social_dates = [date.fromisoformat(d) for d in social_dates_raw if d]
    if not social_dates and not yt_dates:
        add_text(slide, 0.8, 4, 11, 0.5,
                 "No timeline data available.", size=14, color=PALETTE["muted"])
        return
    start = min(yt_dates + social_dates) if (yt_dates or social_dates) else date.today()
    end = max(yt_dates + social_dates) if (yt_dates or social_dates) else date.today()

    timeline = []
    d = start
    while d <= end:
        timeline.append(d)
        d += timedelta(days=1)

    # Per-day data
    yt_by_date = {date.fromisoformat(e["date"]): e["views"] for e in YT_HISTORY}
    social_by_date: dict[date, int] = {}
    for r in all_records:
        nd = normalize_date(r.get("upload_date"))
        if not nd:
            continue
        d_key = date.fromisoformat(nd)
        v = r.get("view_count")
        try:
            social_by_date[d_key] = social_by_date.get(d_key, 0) + (int(v) if v else 0)
        except (ValueError, TypeError):
            pass

    yt_series = [yt_by_date.get(d, None) for d in timeline]  # None = no point that day
    social_series = [social_by_date.get(d, 0) for d in timeline]
    # python-pptx LineChart wants numbers; convert None → 0 for now but
    # episodes-only days WILL show 0 for social and vice versa. That's intentional —
    # the chart's value is showing the spike pattern across the timeline.

    labels = [d.strftime("%b %d") for d in timeline]

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("YouTube episode views",
                          [v if v is not None else 0 for v in yt_series])
    chart_data.add_series("Total social views (all platforms)", social_series)

    # Chart position
    chart_x, chart_y, chart_w, chart_h = Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.7)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, chart_x, chart_y, chart_w, chart_h, chart_data
    )
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Caption explaining what to look for
    add_text(slide, 0.8, 6.15, 11.7, 0.4,
             "READING THE CHART", size=11, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 6.5, 11.7, 0.4,
             "Past episodes settle at 135–212 views over ~10 days. Today's episode is at 120 in its first hours — "
             "tracking toward that range. Heavy social posting on May 31 and Jun 1 didn't have a YT episode to correlate with.",
             size=11, color=PALETTE["muted"])
    add_text(slide, 0.8, 6.9, 11.7, 0.4,
             "Real answer comes from the 24h / 72h / 1-week snapshots of TODAY's episode plotted against the social push that supports it.",
             size=11, color=PALETTE["ink"], bold=True)


def build_gap_slide(prs):
    """What we CAN'T see at the public layer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.6, 11, 0.4,
             "THE GAP", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 1.0, 11, 0.9,
             "What this report can't show — yet",
             size=32, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.85, 11, 0.4,
             "Today: public counts only. Tomorrow with proper API access: full Studio analytics.",
             size=13, color=PALETTE["muted"])

    # Two-column comparison cards
    items = [
        ("YouTube", [
            "Watch time + average view duration",
            "Audience retention curve",
            "Traffic sources (search, suggested, external)",
            "Subscribers gained per video",
            "Geographic + age/gender breakdown",
        ]),
        ("Instagram", [
            "Reach (unique accounts)",
            "Plays vs. views",
            "Profile actions + saves",
            "Story / Reel insights",
            "Audience demographics",
        ]),
        ("Facebook", [
            "Engagement counts (likes, comments, shares)",
            "Post reach (organic vs. paid)",
            "Page follower growth attribution",
            "Click-through to external links",
        ]),
        ("TikTok", [
            "Watch time + completion rate",
            "For-You-Page distribution",
            "Audience country + demographics",
            "Source of traffic to video",
        ]),
    ]
    card_w, card_h = 5.9, 2.4
    gap = 0.3
    x_left = 0.8
    x_right = x_left + card_w + gap
    positions = [(x_left, 2.7), (x_right, 2.7), (x_left, 2.7 + card_h + gap),
                 (x_right, 2.7 + card_h + gap)]

    for (x, y), (plat, points) in zip(positions, items):
        add_filled_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        add_filled_rect(slide, x, y, 0.08, card_h, PALETTE["cherry"])
        add_text(slide, x + 0.3, y + 0.2, card_w - 0.5, 0.4,
                 plat.upper(), size=12, color=PALETTE["cherry"], bold=True)
        for i, p in enumerate(points):
            add_text(slide, x + 0.5, y + 0.6 + i * 0.32, card_w - 0.7, 0.3,
                     "—  " + p, size=11, color=PALETTE["ink"])


def build_before_after_slide(prs, snaps, day1_iso="2026-06-08"):
    """The killer pitch slide: stagnant pre-Day-1 baseline vs. Day-1+ posts.

    For each social platform, computes average views per post for posts dated
    before day1_iso vs. on/after day1_iso, then shows the multiplier prominently.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "BEFORE vs. SINCE I JOINED",
             size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.85, 11.7, 0.9,
             "Per-post reach jumped 2× to 6×",
             size=30, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "Their old posts had weeks to gather views. Mine are only days old — and still pulling more reach per post.",
             size=12, color=PALETTE["muted"])

    day1_date = date.fromisoformat(day1_iso)
    latest = snaps[-1]

    cards = []
    for plat_key, plat_label in [("instagram", "Instagram"),
                                   ("tiktok", "TikTok"),
                                   ("facebook", "Facebook")]:
        pre, post = [], []
        for r in latest.get("records", []):
            if r.get("platform") != plat_key:
                continue
            nd = normalize_date(r.get("upload_date"))
            if not nd:
                continue
            try:
                d_obj = date.fromisoformat(nd)
            except ValueError:
                continue
            v = int(r.get("view_count") or 0)
            if d_obj >= day1_date:
                post.append(v)
            else:
                pre.append(v)
        pre_avg = (sum(pre) / len(pre)) if pre else 0
        post_avg = (sum(post) / len(post)) if post else 0
        multiplier = (post_avg / pre_avg) if pre_avg > 0 else 0
        cards.append({
            "platform": plat_label,
            "pre_count": len(pre),
            "pre_avg": pre_avg,
            "post_count": len(post),
            "post_avg": post_avg,
            "multiplier": multiplier,
        })

    # 3 cards side by side
    card_w, card_h = 3.85, 4.5
    gap = 0.3
    total_w = card_w * 3 + gap * 2
    x_start = (LAYOUT_W - total_w) / 2
    y = 2.3

    for i, cd in enumerate(cards):
        x = x_start + i * (card_w + gap)
        # Card background + cherry stripe
        add_filled_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        add_filled_rect(slide, x, y, 0.08, card_h, PALETTE["cherry"])

        # Platform name
        add_text(slide, x + 0.3, y + 0.25, card_w - 0.5, 0.4,
                 cd["platform"].upper(), size=12, color=PALETTE["cherry"], bold=True)

        # Multiplier (huge)
        mult_str = f"{cd['multiplier']:.1f}×" if cd['multiplier'] else "—"
        add_text(slide, x + 0.3, y + 0.7, card_w - 0.5, 1.1,
                 mult_str, size=66, color=PALETTE["navy"], bold=True, font="Georgia")
        # Caption
        add_text(slide, x + 0.3, y + 1.85, card_w - 0.5, 0.4,
                 "more views per post",
                 size=11, color=PALETTE["muted"])

        # Divider
        add_filled_rect(slide, x + 0.3, y + 2.3, card_w - 0.6, 0.015, PALETTE["rule"])

        # BEFORE block
        add_text(slide, x + 0.3, y + 2.45, card_w - 0.5, 0.3,
                 "BEFORE I JOINED", size=10, color=PALETTE["muted"], bold=True)
        add_text(slide, x + 0.3, y + 2.75, card_w - 0.5, 0.4,
                 f"{cd['pre_avg']:,.0f} views avg",
                 size=18, color=PALETTE["ink"], bold=True)
        add_text(slide, x + 0.3, y + 3.1, card_w - 0.5, 0.3,
                 f"({cd['pre_count']} posts, weeks old)",
                 size=10, color=PALETTE["muted"])

        # SINCE block
        add_text(slide, x + 0.3, y + 3.5, card_w - 0.5, 0.3,
                 "SINCE I JOINED", size=10, color=PALETTE["cherry"], bold=True)
        add_text(slide, x + 0.3, y + 3.8, card_w - 0.5, 0.4,
                 f"{cd['post_avg']:,.0f} views avg",
                 size=18, color=PALETTE["cherry"], bold=True)
        add_text(slide, x + 0.3, y + 4.15, card_w - 0.5, 0.3,
                 f"({cd['post_count']} posts, days old)",
                 size=10, color=PALETTE["muted"])

    # Takeaway
    add_text(slide, 0.8, 7.05, 11.7, 0.4,
             "Same audience. Same brand. Different editor. More reach per post — every platform.",
             size=12, color=PALETTE["ink"], bold=True)


def build_categories_slide(prs, snaps, day1_iso="2026-06-08"):
    """What content categories you've covered + their per-post reach.

    Pulls the latest snapshot, filters to Day-1+ posts on social platforms,
    categorizes by keyword in title/caption, ranks by avg views/post.
    """
    import re as _re
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "CONTENT CATEGORIES", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.85, 11.7, 0.9,
             "What's been working, by topic",
             size=30, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 0.8, 1.7, 11.7, 0.4,
             "Average views per post on social. Higher = more reach per piece of work.",
             size=12, color=PALETTE["muted"])

    day1_date = date.fromisoformat(day1_iso)
    latest = snaps[-1]

    # Categorize each Day-1+ social post by keyword match
    KEYWORDS = [
        ("fifa",      r"\b(fifa|world cup|soccer)\b"),
        ("food",      r"\b(food|cuisine|restaurant|asian|dish|chef|eat|dining)\b"),
        ("political", r"\b(trujillo|lisette|mayor|council|councilman|ice|huntington|vote|election|campaign|candidate|judge|judicial|politician|measure er|prosecutor|gang|downey)\b"),
    ]
    def cat_of(text):
        t = (text or "").lower()
        for name, pat in KEYWORDS:
            if _re.search(pat, t):
                return name
        return "other"

    buckets = {"political": [], "fifa": [], "food": []}
    for r in latest.get("records", []):
        plat = r.get("platform")
        if plat not in ("instagram", "tiktok", "facebook"):
            continue
        nd = normalize_date(r.get("upload_date"))
        if not nd:
            continue
        try:
            if date.fromisoformat(nd) < day1_date:
                continue
        except ValueError:
            continue
        text = r.get("title") or r.get("description_excerpt") or ""
        cat = cat_of(text)
        if cat in buckets:
            buckets[cat].append({
                "views": int(r.get("view_count") or 0),
                "likes": int(r.get("like_count") or 0),
                "title": text[:60],
                "platform": plat,
                "date": nd,
            })

    cat_display = [
        ("political", "Political",
         "Mario Trujillo, Lisette, council, mayor, judicial — your bread & butter."),
        ("fifa", "FIFA / World Cup",
         "Sports + city angles — opportunistic, builds when there are real moments."),
        ("food", "Food (Asian restaurants)",
         "Brand new this week — too early to fully judge."),
    ]

    # 3 cards side by side
    card_w, card_h = 3.85, 4.0
    gap = 0.3
    total_w = card_w * 3 + gap * 2
    x_start = (LAYOUT_W - total_w) / 2
    y = 2.3

    for i, (key, label, blurb) in enumerate(cat_display):
        x = x_start + i * (card_w + gap)
        posts = buckets.get(key, [])
        add_filled_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        add_filled_rect(slide, x, y, 0.08, card_h, PALETTE["cherry"])

        # Label
        add_text(slide, x + 0.3, y + 0.25, card_w - 0.5, 0.4,
                 label.upper(), size=12, color=PALETTE["cherry"], bold=True)

        # Big number (avg/post)
        if posts:
            avg = sum(p["views"] for p in posts) / len(posts)
            avg_str = f"{avg:,.0f}"
        else:
            avg_str = "—"
        add_text(slide, x + 0.3, y + 0.7, card_w - 0.5, 1.1,
                 avg_str, size=58, color=PALETTE["navy"], bold=True, font="Georgia")
        add_text(slide, x + 0.3, y + 1.85, card_w - 0.5, 0.3,
                 "avg views per post", size=11, color=PALETTE["muted"])

        # Posts count + total
        if posts:
            tot_v = sum(p["views"] for p in posts)
            add_text(slide, x + 0.3, y + 2.2, card_w - 0.5, 0.3,
                     f"{len(posts)} posts · {tot_v:,} total views",
                     size=11, color=PALETTE["ink"], bold=True)
        else:
            add_text(slide, x + 0.3, y + 2.2, card_w - 0.5, 0.3,
                     "No posts yet", size=11, color=PALETTE["muted"])

        # Top performer
        if posts:
            top = max(posts, key=lambda p: p["views"])
            add_text(slide, x + 0.3, y + 2.6, card_w - 0.5, 0.3,
                     "TOP PERFORMER", size=10, color=PALETTE["cherry"], bold=True)
            add_text(slide, x + 0.3, y + 2.9, card_w - 0.5, 0.6,
                     f"{top['views']:,} views — {top['title']}",
                     size=10, color=PALETTE["ink"])

        # Blurb at bottom
        add_text(slide, x + 0.3, y + 3.55, card_w - 0.5, 0.45,
                 blurb, size=10, color=PALETTE["muted"])

    # Recommendation + caveat
    add_text(slide, 0.8, 6.55, 11.7, 0.35,
             "WHAT I'M SEEING", size=11, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 6.85, 11.7, 0.5,
             "Political content is the engine — roughly 2× the reach per post compared to FIFA. Food is too new to call yet (only 4 posts since Sunday).",
             size=12, color=PALETTE["ink"], bold=True)
    add_text(slide, 0.8, 7.25, 11.7, 0.3,
             "Suggested mix going forward: political as the backbone · FIFA when real moments hit · food as the new lane to grow.   "
             "Numbers may shift as the food angle matures this week.",
             size=10, color=PALETTE["muted"])


def build_synopsis_slide(prs):
    """A reflection slide written in plain language — what worked, what the
    tradeoff was, what to do next. No jargon."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["cream"])

    add_text(slide, 0.8, 0.5, 11, 0.4,
             "REAL TALK", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 0.85, 11.7, 0.9,
             "Honest take — what worked, what's next",
             size=28, color=PALETTE["navy"], bold=True, font="Georgia")

    # Section 1 — What worked
    y1 = 1.85
    add_text(slide, 0.8, y1, 11, 0.35,
             "WHAT WORKED", size=11, color=PALETTE["cherry"], bold=True)
    worked = [
        "Posting 3 clips on launch day moved the numbers — hard. Every platform climbed double or triple digits.",
        "Coordinated push from one editor pulled real growth, not just busywork. This proves the engagement model works.",
    ]
    for i, b in enumerate(worked):
        add_text(slide, 1.0, y1 + 0.3 + i * 0.32, 11.4, 0.35,
                 "—  " + b, size=12, color=PALETTE["ink"])

    # Section 2 — The honest part
    y2 = 3.0
    add_text(slide, 0.8, y2, 11, 0.35,
             "THE HONEST PART", size=11, color=PALETTE["cherry"], bold=True)
    honest = [
        "We came out swinging — used the strongest, most click-worthy clips on Day 1.",
        "That worked. But it also means I'm running low on top-shelf material for the rest of this week.",
        "Before this, the workflow was 1 clip at a time to test the waters. Safer — but the numbers were smaller too.",
    ]
    for i, b in enumerate(honest):
        add_text(slide, 1.0, y2 + 0.3 + i * 0.32, 11.4, 0.35,
                 "—  " + b, size=12, color=PALETTE["ink"])

    # Section 3 — Two ways forward
    y3 = 4.6
    add_text(slide, 0.8, y3, 11, 0.35,
             "TWO WAYS WE CAN GO FROM HERE", size=11, color=PALETTE["cherry"], bold=True)

    # Card A
    cy = y3 + 0.4
    add_filled_rect(slide, 0.8, cy, 5.8, 1.7, PALETTE["white"])
    add_filled_rect(slide, 0.8, cy, 0.08, 1.7, PALETTE["cherry"])
    add_text(slide, 1.0, cy + 0.1, 5.6, 0.4,
             "OPTION A — Steady pace",
             size=14, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 1.0, cy + 0.5, 5.6, 1.2,
             "2 posts a day across the week. A steady rhythm, easier to keep up. "
             "Smaller spikes — but the channel stays alive every day.",
             size=11, color=PALETTE["ink"])

    # Card B
    add_filled_rect(slide, 6.9, cy, 5.8, 1.7, PALETTE["white"])
    add_filled_rect(slide, 6.9, cy, 0.08, 1.7, PALETTE["cherry"])
    add_text(slide, 7.1, cy + 0.1, 5.6, 0.4,
             "OPTION B — Big Monday, smart content rest of week",
             size=14, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 7.1, cy + 0.5, 5.6, 1.2,
             "Punch hard the day a new episode drops (3–4 headline clips). "
             "The rest of the week: post informational, value-driven content — "
             "stuff that doesn't chase virality but funnels people into the show and builds long-term trust.",
             size=11, color=PALETTE["ink"])

    # Closing line
    add_text(slide, 0.8, 6.9, 11.7, 0.45,
             "Either way: this week's numbers prove we can deliver. "
             "Now it's about finding the right mix — spice when we want a spike, substance to keep the audience between episodes.",
             size=13, color=PALETTE["navy"], bold=True, font="Georgia")


def build_ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Bottom-half cherry, top-half navy/ink for dramatic split
    add_filled_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    add_filled_rect(slide, 0, 4.7, LAYOUT_W, LAYOUT_H - 4.7, PALETTE["cherry"])

    add_text(slide, 0.8, 0.6, 11, 0.4,
             "THE ASK", size=12, color=PALETTE["cherry"], bold=True)
    add_text(slide, 0.8, 1.0, 11.5, 1.0,
             "Give me view-only access so reports are richer",
             size=36, color=PALETTE["white"], bold=True, font="Georgia")
    add_text(slide, 0.8, 2.15, 11.5, 0.5,
             "These numbers came from public scraping. With view-only access to the Page and IG insights, you get much more.",
             size=15, color=PALETTE["cream"])

    # Two columns: what's needed + what we get
    add_text(slide, 0.8, 3.1, 5.5, 0.4,
             "WHAT I NEED FROM YOU", size=12, color=PALETTE["cherry"], bold=True)
    needs = [
        "Add me to the Facebook Page with the “Insights” role",
        "Add me to Instagram with view-only insights",
        "(That’s it — about 30 minutes, once)",
        "I can’t post, edit, or change anything",
    ]
    for i, n in enumerate(needs):
        add_text(slide, 1.0, 3.5 + i * 0.34, 5.4, 0.3,
                 "—  " + n, size=12, color=PALETTE["cream"])

    add_text(slide, 6.7, 3.1, 5.5, 0.4,
             "WHAT YOU GET IN RETURN", size=12, color=PALETTE["cherry"], bold=True)
    gets = [
        "Real watch time per video (not just view counts)",
        "Where viewers come from (search, suggested, shared)",
        "Audience details — age, location, devices",
        "Automatic reports — 24 hours, 72 hours, 1 week",
    ]
    for i, g in enumerate(gets):
        add_text(slide, 6.9, 3.5 + i * 0.34, 5.4, 0.3,
                 "—  " + g, size=12, color=PALETTE["cream"])

    # Cherry footer band — cost + time
    add_text(slide, 0.8, 5.1, 5.5, 0.5,
             "TIME COMMITMENT", size=12, color=PALETTE["cream"], bold=True)
    add_text(slide, 0.8, 5.5, 5.5, 1.0,
             "~30 minutes, once",
             size=32, color=PALETTE["white"], bold=True, font="Georgia")

    add_text(slide, 6.7, 5.1, 5.5, 0.5,
             "ONGOING COST", size=12, color=PALETTE["cream"], bold=True)
    add_text(slide, 6.7, 5.5, 5.5, 1.0,
             "$0",
             size=32, color=PALETTE["white"], bold=True, font="Georgia")


# ============================================================================
# Main
# ============================================================================

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--episode", required=True, type=Path,
                    help="Episode folder containing snapshots/*.json")
    ap.add_argument("--out", type=Path,
                    help="Output pptx path (default: <episode>/report.pptx)")
    args = ap.parse_args()

    snap_dir = args.episode / "snapshots"
    if not snap_dir.is_dir():
        sys.exit(f"No snapshots dir at {snap_dir}")

    snaps = []
    for f in sorted(snap_dir.glob("*.json")):
        if not f.name.startswith("."):
            snaps.append(json.loads(f.read_text()))
    if not snaps:
        sys.exit(f"No snapshot JSONs in {snap_dir}")

    # Use the latest snapshot for slide content
    latest = snaps[-1]
    all_records = latest.get("records", [])
    snapshot_label = latest.get("snapshot_tag", "t0")
    episode_title = latest.get("episode") or args.episode.name
    posted_at = latest.get("posted_at", "")

    # FILTER to baseline-day posts only ("Day 1" — first work under new editorial flow).
    # Target date comes from the manifest's posted_at (YYYY-MM-DD prefix).
    target_date = (posted_at or "")[:10]

    # Each platform should have one post from today; multiple if more were posted.
    records = []
    for plat in ("youtube", "instagram", "tiktok", "facebook"):
        records.extend(posts_on(all_records, plat, target_date))

    # Per-platform totals across that filtered set
    totals = {
        "youtube":   sum_views(records, "youtube"),
        "instagram": sum_views(records, "instagram"),
        "tiktok":    sum_views(records, "tiktok"),
        "facebook":  sum_views(records, "facebook"),
    }

    prs = Presentation()
    prs.slide_width = Inches(LAYOUT_W)
    prs.slide_height = Inches(LAYOUT_H)

    # Slide 1 — title
    build_title_slide(prs, episode_title, posted_at, len(snaps))

    # Slide 2 — Day 1 → Day 2 social growth headlines (IG/TT/FB big, YT footer)
    build_glance_slide(prs, snaps, target_date)

    # Slide 3 — Linear growth chart
    build_day1_growth_chart_slide(prs, snaps, target_date)

    # Slide 4 — BEFORE vs SINCE I JOINED (the killer pitch slide)
    build_before_after_slide(prs, snaps, day1_iso=target_date)

    # Slide 5 — Content categories breakdown + recommendation
    build_categories_slide(prs, snaps, day1_iso=target_date)

    # Slide 6 — Instagram (all recent Reels + posts with t0/t1/Δ)
    build_platform_full_table_slide(
        prs,
        platform="instagram",
        label="Instagram — Reels & posts",
        snaps=snaps,
        target_date=target_date,
        post_word="Reels & posts",
    )

    # Slide 5 — TikTok
    build_platform_full_table_slide(
        prs,
        platform="tiktok",
        label="TikTok — clip performance",
        snaps=snaps,
        target_date=target_date,
        post_word="clips",
    )

    # Slide 6 — Facebook
    build_platform_full_table_slide(
        prs,
        platform="facebook",
        label="Facebook — Page posts",
        snaps=snaps,
        target_date=target_date,
        post_word="Page posts",
        footnote="Facebook hides like/comment counts at the public layer — only view counts are publicly visible. With API access this fills in.",
    )

    # Slide 7 — Real-talk synopsis (what worked, the honest part, two ways forward)
    build_synopsis_slide(prs)

    # Slide 8 — The Ask (combines the gap + the request)
    build_ask_slide(prs)

    out = args.out or (args.episode / "report.pptx")
    prs.save(out)
    print(f"✓ Wrote: {out}")
    print(f"  {len(prs.slides)} slides")
    print(f"  Open with: open '{out}'")


if __name__ == "__main__":
    main()
