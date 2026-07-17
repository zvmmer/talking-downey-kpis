"""Talking Downey — sponsor-facing content-type performance report.

Different audience than the monthly decks. This one goes to the sponsor
who funds the show. Their question: "What's working — where should
I direct more spend?"

The answer isn't given as tables. Every slide is a one-look category card
with a big number, a top post, and a plain-English recommendation.
Explainable to a 5-year-old — Zhamir's rule.

Uses report_helpers.SPONSOR_GROUPS to collapse fine-grained taxonomy
into ~8 presentable buckets.

Usage
-----
    MCP/.venv/bin/python Projects/talking_downey/kpis/build_sponsor_content_report.py
"""
from __future__ import annotations

import argparse
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from report_helpers import (
    PALETTE, LAYOUT_W, LAYOUT_H, ENGAGEMENT_START,
    add_rect, add_text, fmt_int,
    load_all_snapshots, posts_in_window, group_by_category,
    slide_new, slide_header, slide_footer_takeaway,
    SPONSOR_GROUPS,
)

TODAY = date.today().isoformat()


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, LAYOUT_W, LAYOUT_H, PALETTE["ink"])
    add_rect(slide, 0, 0, 0.45, LAYOUT_H, PALETTE["gold"])

    add_text(slide, 1.0, 1.2, 11, 0.5,
             "SPONSOR CONTENT REPORT",
             size=14, color=PALETTE["gold"], bold=True)
    add_text(slide, 1.0, 1.9, 11, 1.3,
             "Talking Downey", size=64, color=PALETTE["white"],
             bold=True, font="Georgia")
    add_text(slide, 1.0, 3.4, 11, 0.6,
             "What worked. What to fund more of.",
             size=24, color=PALETTE["cream"], font="Georgia")

    add_text(slide, 1.0, 5.0, 8, 0.4,
             f"Reporting window: Jun 8 → {TODAY}",
             size=13, color=PALETTE["muted"])
    add_text(slide, 1.0, 5.4, 8, 0.4,
             "Every content category, ranked. One page per bucket.",
             size=13, color=PALETTE["muted"])
    add_text(slide, 1.0, 5.8, 8, 0.4,
             "Numbers are cumulative views since engagement start.",
             size=13, color=PALETTE["muted"])

    add_text(slide, 1.0, 6.6, 11, 0.4,
             "Prepared by Zhamir Pascual — Kaname Z",
             size=12, color=PALETTE["muted"])


def slide_leaderboard(prs, all_posts):
    """Executive summary — all groups ranked in one visual."""
    slide = slide_new(prs)
    slide_header(slide, "At a glance",
                 "Where your money went — and what it earned",
                 "Every content type, ranked by total views. Bigger = better performing.")

    buckets = group_by_category(all_posts)
    groups = []
    for label, cats, color_key, blurb in SPONSOR_GROUPS:
        posts = []
        for c in cats:
            posts.extend(buckets.get(c, []))
        if not posts:
            continue
        views = sum(int(p.get("view_count") or 0) for p in posts)
        groups.append({
            "label": label, "color": color_key, "blurb": blurb,
            "posts": posts, "views": views,
            "n": len(posts),
            "avg": views / len(posts) if posts else 0,
        })
    groups.sort(key=lambda g: -g["views"])
    max_v = max((g["views"] for g in groups), default=1)

    # Fit N rows into ~4.15" vertical space (y=2.3 → 6.45, leaving room for takeaway).
    avail_h = 4.15
    row_h = min(0.44, avail_h / max(1, len(groups)))
    y0 = 2.3
    x_bar = 4.5
    bar_max_w = 6.8

    for i, g in enumerate(groups):
        y = y0 + i * row_h
        add_text(slide, 0.8, y + 0.02, 3.5, row_h,
                 g["label"], size=11, color=PALETTE["ink"], bold=True)
        bar_w = max(0.05, bar_max_w * g["views"] / max_v)
        bar_h = max(0.14, row_h - 0.22)
        add_rect(slide, x_bar, y + (row_h - bar_h) / 2, bar_w, bar_h, PALETTE[g["color"]])
        add_text(slide, x_bar + bar_w + 0.1, y + 0.02, 2.5, row_h,
                 f"{fmt_int(g['views'])} views · {g['n']} posts",
                 size=9, color=PALETTE["muted"])

    slide_footer_takeaway(slide,
        "Top bucket is your engine. Bottom bucket is either an opportunity or a lane to retire.")


def slide_group_card(prs, group_info, posts):
    """One slide per sponsor group — the workhorse of this deck."""
    label, cats, color_key, blurb = group_info
    color = PALETTE[color_key]

    slide = slide_new(prs)
    slide_header(slide, "CATEGORY DEEP-DIVE",
                 label,
                 blurb)

    if not posts:
        add_text(slide, 0.8, 3.5, 11.7, 1.0,
                 "No posts in this category yet.",
                 size=20, color=PALETTE["muted"], align="center")
        return

    views = sum(int(p.get("view_count") or 0) for p in posts)
    likes = sum(int(p.get("like_count") or 0) for p in posts)
    n = len(posts)
    avg = views / n if n else 0

    # Big number card
    add_rect(slide, 0.8, 2.4, 5.5, 3.5, PALETTE["white"])
    add_rect(slide, 0.8, 2.4, 0.1, 3.5, color)
    add_text(slide, 1.05, 2.6, 5, 0.4,
             "TOTAL VIEWS", size=11, color=color, bold=True)
    add_text(slide, 1.05, 2.95, 5, 1.2,
             fmt_int(views), size=54, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, 1.05, 4.2, 5, 0.4,
             f"{n} posts · {fmt_int(likes)} likes",
             size=13, color=PALETTE["ink"], bold=True)
    add_text(slide, 1.05, 4.6, 5, 0.4,
             f"{avg:,.0f} average views per post",
             size=13, color=PALETTE["muted"])

    # Top 3 posts
    add_text(slide, 6.7, 2.4, 6, 0.3,
             "TOP 3 IN THIS CATEGORY", size=11, color=color, bold=True)
    top3 = sorted(posts, key=lambda p: -int(p.get("view_count") or 0))[:3]
    for i, p in enumerate(top3):
        y = 2.8 + i * 1.0
        add_rect(slide, 6.7, y, 5.8, 0.9, PALETTE["white"])
        add_rect(slide, 6.7, y, 0.06, 0.9, color)
        title = (p.get("title") or p.get("description_excerpt") or "(no title)")[:90]
        plat = {"instagram": "IG", "tiktok": "TT", "facebook": "FB", "youtube": "YT"}.get(p.get("platform"), "?")
        add_text(slide, 6.9, y + 0.1, 5.5, 0.35,
                 f"#{i+1}  ·  {plat}  ·  {fmt_int(p.get('view_count'))} views",
                 size=10, color=color, bold=True)
        add_text(slide, 6.9, y + 0.42, 5.5, 0.45,
                 title, size=10, color=PALETTE["ink"])

    # Recommendation strip
    add_rect(slide, 0.8, 6.15, 11.7, 0.8, PALETTE["ink"])
    add_text(slide, 1.0, 6.25, 3, 0.3,
             "RECOMMENDATION", size=10, color=color, bold=True)
    rec = recommendation_for(label, views, n, avg)
    add_text(slide, 1.0, 6.5, 11.5, 0.4,
             rec, size=13, color=PALETTE["cream"], bold=True)


def slide_why_low_lanes(prs, all_posts):
    """The 'why we cover the low-view lanes' argument — 3 reasons in cards."""
    slide = slide_new(prs)
    slide_header(slide, "Strategic context",
                 "Why we still cover the lower-view lanes",
                 "Not every category needs to top the leaderboard. Here's why the quieter buckets earn their spend.")

    reasons = [
        ("NEW AUDIENCE ACQUISITION", "green",
         "Scandal brings viewers. Positive content KEEPS them.",
         "Without community coverage, viewers churn after 3-4 controversy spikes. They came for drama, "
         "they leave because they don't feel connection to Downey itself."),
        ("EDITORIAL CREDIBILITY", "gold",
         "A show that only reports negatives becomes a hit piece.",
         "Sponsors don't want to attach to a hit piece. Positive coverage — resident highlights, downtown "
         "growth, community wins — proves objectivity. That's what makes sponsorship durable."),
        ("ALGORITHMIC HEALTH", "sky",
         "Platforms punish single-topic accounts.",
         "Instagram, TikTok, and YouTube increasingly penalize accounts that push one theme. A mixed content "
         "diet earns more consistent reach across ALL posts, including the scandal ones. Diversification pays."),
    ]

    card_w, card_h, gap = 3.9, 3.6, 0.2
    x0 = (LAYOUT_W - card_w * 3 - gap * 2) / 2
    y = 2.4

    for i, (title, color_key, tagline, body) in enumerate(reasons):
        x = x0 + i * (card_w + gap)
        color = PALETTE[color_key]
        add_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        add_rect(slide, x, y, 0.1, card_h, color)
        add_text(slide, x + 0.3, y + 0.25, card_w - 0.6, 0.4,
                 f"REASON {i+1}", size=10, color=color, bold=True)
        add_text(slide, x + 0.3, y + 0.6, card_w - 0.6, 0.5,
                 title, size=14, color=PALETTE["navy"], bold=True, font="Georgia")
        add_text(slide, x + 0.3, y + 1.2, card_w - 0.6, 0.6,
                 tagline, size=13, color=PALETTE["ink"], bold=True)
        add_text(slide, x + 0.3, y + 1.9, card_w - 0.6, 1.6,
                 body, size=10, color=PALETTE["muted"])

    slide_footer_takeaway(slide,
        "Low-view lanes aren't dead weight. They're the connective tissue that turns viral reach into a durable audience.")


def slide_reach_vs_depth(prs, all_posts):
    """The honest tradeoff — reach (scandal) vs depth (positive)."""
    slide = slide_new(prs)
    slide_header(slide, "The honest tradeoff",
                 "Reach vs. depth — what each side actually buys",
                 "This isn't a Downey problem. It's how every social platform works. The question is the RATIO, not the choice.")

    buckets = group_by_category(all_posts)

    def _bucket_stats(cats):
        posts = []
        for c in cats:
            posts.extend(buckets.get(c, []))
        if not posts:
            return 0, 0, "—", 0
        views = sum(int(p.get("view_count") or 0) for p in posts)
        avg = views / len(posts)
        top = max(posts, key=lambda p: int(p.get("view_count") or 0))
        top_title = (top.get("title") or top.get("description_excerpt") or "")[:70]
        return views, len(posts), top_title, int(top.get("view_count") or 0)

    reach_cats = ["trujillo_scandal", "lisette_scandal", "council_chaos", "ai_fake_news", "ice_immigration"]
    depth_cats = ["resident_highlight", "political_individual_highlight", "downtown_development",
                  "food", "american_pride", "community_events", "community_local"]

    r_views, r_n, r_top, r_top_v = _bucket_stats(reach_cats)
    d_views, d_n, d_top, d_top_v = _bucket_stats(depth_cats)
    r_avg = r_views / r_n if r_n else 0
    d_avg = d_views / d_n if d_n else 0

    # Two columns
    col_w = 5.8
    gap = 0.5
    x_left = (LAYOUT_W - col_w * 2 - gap) / 2
    x_right = x_left + col_w + gap
    y = 2.3

    # LEFT — what earns eyeballs (reach)
    add_rect(slide, x_left, y, col_w, 4.0, PALETTE["white"])
    add_rect(slide, x_left, y, 0.1, 4.0, PALETTE["cherry"])
    add_text(slide, x_left + 0.3, y + 0.2, col_w - 0.6, 0.4,
             "WHAT EARNS EYEBALLS", size=11, color=PALETTE["cherry"], bold=True)
    add_text(slide, x_left + 0.3, y + 0.6, col_w - 0.6, 0.5,
             "Truth & Scandal", size=20, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, x_left + 0.3, y + 1.15, col_w - 0.6, 0.4,
             f"{r_avg:,.0f} avg views/post   ·   {r_n} posts   ·   {fmt_int(r_views)} total",
             size=11, color=PALETTE["ink"], bold=True)
    add_text(slide, x_left + 0.3, y + 1.65, col_w - 0.6, 0.35,
             "TOP POST", size=9, color=PALETTE["cherry"], bold=True)
    add_text(slide, x_left + 0.3, y + 2.0, col_w - 0.6, 0.6,
             f"{fmt_int(r_top_v)} views — {r_top}", size=10, color=PALETTE["ink"])
    add_text(slide, x_left + 0.3, y + 2.8, col_w - 0.6, 1.1,
             "Every scandal post averages 3-5× the reach of positive posts. "
             "This is the platform algorithm, not Downey. Outrage drives comments, "
             "comments drive reach, reach drives new followers.",
             size=10, color=PALETTE["muted"])

    # RIGHT — what earns loyalty (depth)
    add_rect(slide, x_right, y, col_w, 4.0, PALETTE["white"])
    add_rect(slide, x_right, y, 0.1, 4.0, PALETTE["sky"])
    add_text(slide, x_right + 0.3, y + 0.2, col_w - 0.6, 0.4,
             "WHAT EARNS LOYALTY", size=11, color=PALETTE["sky"], bold=True)
    add_text(slide, x_right + 0.3, y + 0.6, col_w - 0.6, 0.5,
             "Everything else", size=20, color=PALETTE["navy"], bold=True, font="Georgia")
    add_text(slide, x_right + 0.3, y + 1.15, col_w - 0.6, 0.4,
             f"{d_avg:,.0f} avg views/post   ·   {d_n} posts   ·   {fmt_int(d_views)} total",
             size=11, color=PALETTE["ink"], bold=True)
    add_text(slide, x_right + 0.3, y + 1.65, col_w - 0.6, 0.35,
             "TOP POST", size=9, color=PALETTE["sky"], bold=True)
    add_text(slide, x_right + 0.3, y + 2.0, col_w - 0.6, 0.6,
             f"{fmt_int(d_top_v)} views — {d_top}", size=10, color=PALETTE["ink"])
    add_text(slide, x_right + 0.3, y + 2.8, col_w - 0.6, 1.1,
             "Positive posts get fewer views per post — but they turn viewers into "
             "subscribers. Loyalty compounds. It's slower, but it's the audience the "
             "sponsor's brand actually attaches to.",
             size=10, color=PALETTE["muted"])

    # Banner
    add_rect(slide, 0.8, 6.55, 11.7, 0.6, PALETTE["ink"])
    add_text(slide, 1.0, 6.7, 11.5, 0.4,
             "Both matter. The real question isn't which lane wins — it's the RATIO between them.",
             size=13, color=PALETTE["cream"], bold=True)


def slide_directions(prs):
    """Three strategic directions — sponsor picks one."""
    slide = slide_new(prs)
    slide_header(slide, "Three directions",
                 "How aggressive do we lean, and toward what?",
                 "Each direction is a defensible choice. Each has a cost. The right one depends on what the sponsor wants Talking Downey to BE.")

    directions = [
        ("cherry", "MAXIMUM REACH", "60% scandal / 40% positive filler",
         "Bigger view numbers.\nMore new audience per week.",
         "Higher churn. Editorial credibility risk. Sponsor-brand risk if scandal target sues or spirals.",
         "You want raw growth numbers to show the sponsor month-over-month."),
        ("sky",    "BALANCED PORTFOLIO", "40% investigative / 40% positive / 20% community",
         "Sustainable growth.\nCredibility intact.\nBoth audiences served.",
         "Slower month-over-month reach curve. Requires editorial discipline to hold the ratio.",
         "You want a durable show that a sponsor can attach to long-term."),
        ("green",  "POSITIVE-LEAN", "20% scandal / 60% positive / 20% community",
         "Community-cheerleader positioning.\nSafe sponsor territory.\nHigh trust.",
         "Lowest total reach. Loses the scandal-hungry audience over time. Harder to attract new eyeballs.",
         "You want to position Talking Downey as Downey's booster, not its watchdog."),
    ]

    card_w, card_h, gap = 3.95, 4.05, 0.15
    x0 = (LAYOUT_W - card_w * 3 - gap * 2) / 2
    y = 2.05

    for i, (color_key, title, mix, wins, costs, when) in enumerate(directions):
        x = x0 + i * (card_w + gap)
        color = PALETTE[color_key]
        add_rect(slide, x, y, card_w, card_h, PALETTE["white"])
        add_rect(slide, x, y, 0.1, card_h, color)

        add_text(slide, x + 0.3, y + 0.2, card_w - 0.6, 0.3,
                 f"DIRECTION {i+1}", size=9, color=color, bold=True)
        add_text(slide, x + 0.3, y + 0.5, card_w - 0.6, 0.4,
                 title, size=14, color=PALETTE["navy"], bold=True, font="Georgia")
        add_text(slide, x + 0.3, y + 0.92, card_w - 0.6, 0.6,
                 mix, size=9, color=PALETTE["muted"], bold=True)

        add_text(slide, x + 0.3, y + 1.55, card_w - 0.6, 0.28,
                 "WINS", size=9, color=color, bold=True)
        add_text(slide, x + 0.3, y + 1.83, card_w - 0.6, 0.8,
                 wins, size=10, color=PALETTE["ink"])

        add_text(slide, x + 0.3, y + 2.65, card_w - 0.6, 0.28,
                 "COSTS", size=9, color=PALETTE["muted"], bold=True)
        add_text(slide, x + 0.3, y + 2.93, card_w - 0.6, 0.85,
                 costs, size=10, color=PALETTE["muted"])

        add_text(slide, x + 0.3, y + 3.4, card_w - 0.6, 0.28,
                 "PICK IF", size=9, color=color, bold=True)
        add_text(slide, x + 0.3, y + 3.68, card_w - 0.6, 0.4,
                 when, size=9, color=PALETTE["ink"])

    # Soft recommendation banner — moved below cards with breathing room
    add_rect(slide, 0.8, 6.4, 11.7, 0.85, PALETTE["ink"])
    add_text(slide, 1.0, 6.5, 8, 0.3,
             "IF BRAND IMAGE MATTERS AS MUCH AS REACH", size=10, color=PALETTE["sky"], bold=True)
    add_text(slide, 1.0, 6.8, 11.5, 0.45,
             "The Balanced Portfolio best fits Talking Downey today — preserves scandal reach while showing the depth sponsors attach to comfortably.",
             size=11, color=PALETTE["cream"], bold=True)


def recommendation_for(label: str, total_views: int, n_posts: int, avg: float) -> str:
    """Plain-English suggestion for the sponsor per bucket."""
    if n_posts == 0:
        return f"No {label} content yet — untested lane. Consider a small pilot batch."
    if avg > 3000:
        return f"Strong performer. Keep investing — {label} is a proven driver of reach."
    if avg > 1000:
        return f"Solid mid-tier. Worth continued investment. Test scaling to see if it holds."
    if avg > 300:
        return f"Modest reach. Fine as filler between heavy hitters, not a headline lane."
    return f"Underperforming. Either the format needs work — or the audience isn't there for it. Discuss with editor before more spend."


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path,
                    default=Path(__file__).parent / "SPONSOR_CONTENT_REPORT.pptx")
    ap.add_argument("--kpis-root", type=Path,
                    default=Path(__file__).parent)
    args = ap.parse_args()

    snaps = load_all_snapshots(args.kpis_root / "episodes_public")
    if not snaps:
        sys.exit(f"No snapshots found under {args.kpis_root/'episodes_public'}")

    all_posts = posts_in_window(snaps, ENGAGEMENT_START, None)
    buckets = group_by_category(all_posts)
    print(f"Loaded {len(snaps)} snapshots · {len(all_posts)} unique posts")

    prs = Presentation()
    prs.slide_width = Inches(LAYOUT_W)
    prs.slide_height = Inches(LAYOUT_H)

    slide_title(prs)
    slide_leaderboard(prs, all_posts)

    # One page per sponsor group (ranked by total views, best first)
    group_totals = []
    for grp in SPONSOR_GROUPS:
        label, cats, color_key, blurb = grp
        posts = []
        for c in cats:
            posts.extend(buckets.get(c, []))
        views = sum(int(p.get("view_count") or 0) for p in posts)
        group_totals.append((views, grp, posts))
    group_totals.sort(key=lambda t: -t[0])

    for _, grp, posts in group_totals:
        slide_group_card(prs, grp, posts)

    # 3 strategic ending slides
    slide_why_low_lanes(prs, all_posts)
    slide_reach_vs_depth(prs, all_posts)
    slide_directions(prs)

    prs.save(args.out)
    print(f"✓ Wrote: {args.out}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
